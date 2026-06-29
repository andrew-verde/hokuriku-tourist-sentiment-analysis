#!/usr/bin/env python3
"""Build NUDGE-Seminar-Slides.pptx - a native, editable PowerPoint of the nudge deck.

Same provenance contract as the HTML deck: every number resolves from an analysis
output file via the SAME getters. We do not re-implement getters or re-type digits.
We import scripts/build_nudge_seminar_slides.py, call deck.load(), and pull RAW values
straight off deck.DATA[...] through the deck's getter factories. If a getter raises,
the build crashes (fail-loud) rather than inventing a number. Figures are the same
script-generated SVGs, rasterized to PNG with cairosvg.

Style: school colors only (white, blue, black). Century Gothic. One section label
top-left. No em dashes. Stats-heavy methods and stats-justified results.

Run:  .venv/bin/python3 scripts/build_nudge_pptx.py
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "NUDGE-Seminar-Slides.pptx"
SCRATCH = Path(
    "/tmp/claude-1000/-home-andrewgreen-Repositories-andrew-verde-hokuriku-tourist-sentiment-analysis/"
    "10c07a34-3a17-4c5a-aa81-7faf2262a082/scratchpad"
)
SCRATCH.mkdir(parents=True, exist_ok=True)
N_SLIDES = 13


# --- import the HTML deck module and reuse its getters / SOURCES / FIGURES -----
def _load_deck():
    spec = importlib.util.spec_from_file_location(
        "nudge_deck", ROOT / "scripts/build_nudge_seminar_slides.py"
    )
    mod = importlib.util.module_from_spec(spec)
    sys.modules["nudge_deck"] = mod
    spec.loader.exec_module(mod)
    mod.load()  # populate DATA / SHA
    return mod


deck = _load_deck()


# --- value resolution: RAW values come ONLY through deck getters --------------
def raw(src_id: str, getter):
    """Resolve a single raw value via a deck getter (fail-loud)."""
    return getter(deck.DATA[src_id])


_SUP = str.maketrans("-0123456789", "⁻⁰¹²³⁴⁵⁶⁷⁸⁹")


def fp(p: float) -> str:
    """'p = ...' display with unicode superscripts for small p."""
    p = float(p)
    if p < 1e-4:
        mant, exp = f"{p:.1e}".split("e")
        return f"p = {mant}×10{str(int(exp)).translate(_SUP)}"
    if p < 0.001:
        return "p < .001"
    return f"p = {p:.3f}".replace("0.", ".")


def fpv(p: float) -> str:
    """Bare p value (no 'p = ' prefix), for table cells."""
    p = float(p)
    if p < 1e-4:
        mant, exp = f"{p:.1e}".split("e")
        return f"{mant}×10{str(int(exp)).translate(_SUP)}"
    if p < 0.001:
        return "< .001"
    return f"{p:.3f}".replace("0.", ".")


def fpct1(x) -> str:
    return f"{float(x) * 100:.1f}%"


def fnum(x) -> str:
    return f"{int(round(float(x))):,}"


def f2(x) -> str:
    return f"{float(x):.2f}"


def ftext(x) -> str:
    return f"{x}"


# ---- pull EVERY displayed number through a deck getter ------------------------
# Scope / data
TOTAL_REVIEWS = fnum(raw("nudge_poi", deck.poi_total_reviews()))
N_POIS = fnum(raw("poi_mfst", deck.manifest_metric("n_pois_total")))
TAGGED_ROWS = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_input_rows")))
LANG_JP = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_language_group_counts", "japanese")))
LANG_EN = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_language_group_counts", "english")))
LANG_CN = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_language_group_counts", "chinese")))
LANG_OTHER = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_language_group_counts", "other_non_english_non_japanese")))
LANG_UNDETECTED = fnum(raw("aspect_mfst", deck.manifest_metric("tagged_language_group_counts", "undetected_or_too_short")))
POI_FUKUI = fnum(raw("poi_mfst", deck.manifest_metric("n_pois_by_prefecture", "Fukui")))
POI_ISHIKAWA = fnum(raw("poi_mfst", deck.manifest_metric("n_pois_by_prefecture", "Ishikawa")))
POI_TOYAMA = fnum(raw("poi_mfst", deck.manifest_metric("n_pois_by_prefecture", "Toyama")))

# Pipeline
MODEL_ROWS = fnum(raw("aspect_mfst", deck.manifest_metric("primary_model_rows")))
LOW_RATING_ROWS = fnum(raw("aspect_mfst", deck.manifest_metric("primary_low_rating_rows")))
LOW_RATING_DEF = ftext(raw("aspect_mfst", deck.manifest_top("filters", "low_rating_definition")))
LOW_RATING_STARS = fnum(raw("aspect_mfst", deck.low_rating_cutoff()))
N_ASPECTS = fnum(raw("nudge_tax", deck.taxonomy_count()))

# Model + guards
MODELS_PRIMARY = fnum(raw("aspect_mfst", deck.manifest_metric("A_primary_models_fit")))
MODELS_TOTAL = fnum(raw("aspect_mfst", deck.manifest_metric("total_models_fit")))
STATUS_OK = fnum(raw("aspect_mfst", deck.manifest_metric("status_counts", "ok")))
STATUS_PREV = fnum(raw("aspect_mfst", deck.manifest_metric("status_counts", "prevalence_only")))
STATUS_SKIP = fnum(raw("aspect_mfst", deck.manifest_metric("status_counts", "skipped")))
MIN_P = fp(raw("aspect_mfst", deck.manifest_metric("min_model_p_value")))
FIRTH_SANITY = f"{float(raw('aspect_mfst', deck.plain_vs_firth_max_logdiff())):.3f}"
FIRTH_SANITY_ASPECTS = fnum(raw("aspect_mfst", deck.plain_vs_firth_aspect_count()))
CAVEAT_GATE = ftext(raw("aspect_mfst", deck.caveat_text(6)))
MIN_RANKING_MENTIONS = fnum(raw("aspect_mfst", deck.manifest_top("filters", "min_pooled_positive_for_ranking")))


def asp(aspect, col):
    return raw("nudge_aspect", deck.aspect_value(aspect, col))


# Aspect opportunity map (prevalence, odds ratio, 95% CI, FDR p)
OPEN_PREV = fpct1(asp("opening_hours_availability", "prevalence"))
OPEN_N = fnum(asp("opening_hours_availability", "n_positive"))
OPEN_OR = f2(asp("opening_hours_availability", "odds_ratio"))
OPEN_CIL = f2(asp("opening_hours_availability", "or_ci_low"))
OPEN_CIH = f2(asp("opening_hours_availability", "or_ci_high"))
OPEN_P = fp(asp("opening_hours_availability", "p_value_bh_fdr"))
OPEN_PV = fpv(asp("opening_hours_availability", "p_value_bh_fdr"))

TIME_PREV = fpct1(asp("itinerary_fit_time_cost", "prevalence"))
TIME_N = fnum(asp("itinerary_fit_time_cost", "n_positive"))
TIME_OR = f2(asp("itinerary_fit_time_cost", "odds_ratio"))
TIME_CIL = f2(asp("itinerary_fit_time_cost", "or_ci_low"))
TIME_CIH = f2(asp("itinerary_fit_time_cost", "or_ci_high"))
TIME_P = fp(asp("itinerary_fit_time_cost", "p_value_bh_fdr"))
TIME_PV = fpv(asp("itinerary_fit_time_cost", "p_value_bh_fdr"))

PRICE_PREV = fpct1(asp("price_value", "prevalence"))
PRICE_OR = f2(asp("price_value", "odds_ratio"))
PRICE_CIL = f2(asp("price_value", "or_ci_low"))
PRICE_CIH = f2(asp("price_value", "or_ci_high"))
PRICE_PV = fpv(asp("price_value", "p_value_bh_fdr"))

CLEAN_PREV = fpct1(asp("cleanliness_comfort", "prevalence"))
CLEAN_OR = f2(asp("cleanliness_comfort", "odds_ratio"))
CLEAN_CIL = f2(asp("cleanliness_comfort", "or_ci_low"))
CLEAN_CIH = f2(asp("cleanliness_comfort", "or_ci_high"))
CLEAN_PV = fpv(asp("cleanliness_comfort", "p_value_bh_fdr"))

ENG_GAP_PREV = fpct1(asp("english_information_gap", "prevalence"))
SIGN_PREV = fpct1(asp("wayfinding_signage", "prevalence"))
SIGN_N = fnum(asp("wayfinding_signage", "n_positive"))
SIGN_OR = f2(asp("wayfinding_signage", "odds_ratio"))
SIGN_CIL = f2(asp("wayfinding_signage", "or_ci_low"))
SIGN_CIH = f2(asp("wayfinding_signage", "or_ci_high"))
SIGN_PV = fpv(asp("wayfinding_signage", "p_value_bh_fdr"))
ACCESS_PREV = fpct1(asp("transport_access", "prevalence"))
BOOKING_PREV = fpct1(asp("booking_ticketing", "prevalence"))

# Chinese-language Xiaohongshu context
CN_ROWS = fnum(raw("cn_mfst", deck.manifest_metric("denominators", "chinese_social_rows")))
CN_XHS_ROWS = fnum(raw("cn_mfst", deck.manifest_metric("denominators", "n_total_xhs_rows")))


def cnt(code, col):
    return raw("cn_topics", deck.cn_topic_value(code, col))


XHS_SCENIC_N = fnum(cnt("scenic_nature", "count"))
XHS_SCENIC_PCT = f"{float(cnt('scenic_nature', 'pct_posts')):.1f}%"
XHS_DINO_N = fnum(cnt("dinosaurs_museums", "count"))
XHS_DINO_PCT = f"{float(cnt('dinosaurs_museums', 'pct_posts')):.1f}%"


def cnv(predictor, col):
    return raw(
        "cn_drivers",
        deck.cn_driver_value(predictor, "sentiment_category=positive", col),
    )


DINO_N = fnum(cnv("dinosaurs_museums", "group_a_n"))
DINO_POS = fnum(cnv("dinosaurs_museums", "group_a_event_count"))
DINO_PCT = fpct1(cnv("dinosaurs_museums", "group_a_event_pct"))
DINO_OTHER_PCT = fpct1(cnv("dinosaurs_museums", "group_b_event_pct"))
DINO_FDR = fp(cnv("dinosaurs_museums", "p_value_bh_fdr"))
SCENIC_N = fnum(cnv("scenic_nature", "group_a_n"))
SCENIC_POS = fnum(cnv("scenic_nature", "group_a_event_count"))
SCENIC_PCT = fpct1(cnv("scenic_nature", "group_a_event_pct"))
SCENIC_OTHER_PCT = fpct1(cnv("scenic_nature", "group_b_event_pct"))
SCENIC_FDR = fp(cnv("scenic_nature", "p_value_bh_fdr"))


def prv(rank, col):
    return raw("solution_priorities", deck.priority_value(rank, col))


PRIORITIES = [
    {
        "rank": fnum(prv(rank, "rank")),
        "name_en": ftext(prv(rank, "solution_label_en")),
        "name_ja": ftext(prv(rank, "solution_label_ja")),
        "impact": ftext(prv(rank, "impact_tier")),
        "ease": ftext(prv(rank, "ease_tier")),
        "summary_en": ftext(prv(rank, "evidence_summary_en")),
        "summary_ja": ftext(prv(rank, "evidence_summary_ja")),
        "test_en": ftext(prv(rank, "intervention_en")),
        "test_ja": ftext(prv(rank, "intervention_ja")),
    }
    for rank in (1, 2, 3)
]

# POI action map
FIX_COUNT = fnum(raw("nudge_poi", deck.poi_sum("is_fix_it")))
FIX_FUKUI = fnum(raw("nudge_poi", deck.poi_sum("is_fix_it", True)))
PROMOTE_COUNT = fnum(raw("nudge_poi", deck.poi_sum("is_promote_it")))
PROMOTE_FUKUI = fnum(raw("nudge_poi", deck.poi_sum("is_promote_it", True)))
PROMOTE_STRICT_COUNT = fnum(raw("nudge_poi", deck.poi_sum("is_promote_it_strict")))
CROWD_COUNT = fnum(raw("nudge_poi", deck.poi_sum("is_crowding_hotspot")))
LOW_VOL = fnum(raw("nudge_poi", deck.poi_metric("low_volume_threshold")))
HIGH_VOL = fnum(raw("nudge_poi", deck.poi_metric("high_volume_threshold")))
LOW_CONFIDENCE = fnum(raw("poi_mfst", deck.manifest_metric("low_confidence_threshold_n_reviews")))
PROMO1 = ftext(raw("nudge_poi", deck.poi_ranked("promote_fukui", 0, "poi_name")))
PROMO1_N = fnum(raw("nudge_poi", deck.poi_ranked("promote_fukui", 0, "n_reviews")))
PROMO1_SHARE = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 0, "positive_share")))
PROMO1_LOW = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 0, "positive_share_ci_low")))
PROMO1_HIGH = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 0, "positive_share_ci_high")))
PROMO2 = ftext(raw("nudge_poi", deck.poi_ranked("promote_fukui", 1, "poi_name")))
PROMO2_N = fnum(raw("nudge_poi", deck.poi_ranked("promote_fukui", 1, "n_reviews")))
PROMO2_SHARE = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 1, "positive_share")))
PROMO2_LOW = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 1, "positive_share_ci_low")))
PROMO2_HIGH = fpct1(raw("nudge_poi", deck.poi_ranked("promote_fukui", 1, "positive_share_ci_high")))

# Discussion caveats (verbatim manifest strings)
CAV0 = ftext(raw("aspect_mfst", deck.caveat_text(0)))
CAV1 = ftext(raw("aspect_mfst", deck.caveat_text(1)))
CAV3 = ftext(raw("aspect_mfst", deck.caveat_text(3)))
CAV8 = ftext(raw("aspect_mfst", deck.caveat_text(8)))


# ---- figures: rasterize each SVG in deck.FIGURES to PNG ----------------------
def rasterize(key: str) -> Path:
    svg_path = deck.FIGURES[key]
    if not svg_path.exists():
        raise SystemExit(f"missing figure: {svg_path}")
    png = SCRATCH / f"nudge_pptx_{key}.png"
    cairosvg.svg2png(url=str(svg_path), write_to=str(png), output_width=1600)
    return png


PNG = {key: rasterize(key) for key in deck.FIGURES}


# --- palette: white / blue / black only (school colors) -----------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x16, 0x34, 0x5E)   # primary: headings, labels, rules, key numbers
BLUE = RGBColor(0x2C, 0x5A, 0xB0)   # accent
INK = RGBColor(0x14, 0x18, 0x1D)    # body text
GREY = RGBColor(0x5A, 0x65, 0x73)   # secondary + Japanese text
BORDER = RGBColor(0xD3, 0xD9, 0xE2)
CARD_BG = RGBColor(0xF4, 0xF6, 0xF8)
TINT = RGBColor(0xEA, 0xF0, 0xF8)   # very light blue box (still blue family)

EN_FONT = "Century Gothic"
HEAD_FONT = "Century Gothic"
JP_FONT = "Yu Gothic"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
MX = Inches(0.6)
CW = Inches(13.333 - 1.2)


# --- low-level helpers --------------------------------------------------------
def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_run(run, size, color, bold=False, italic=False, font=EN_FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def add_para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    return p


def en_jp(tf, en, jp, en_size=18, jp_size=14, en_color=INK, bold=False, first=False,
          space_after=10, bullet=False):
    """Add an EN paragraph then a JP paragraph directly below (lighter, JP font)."""
    pe = add_para(tf, first=first)
    pe.space_after = Pt(1)
    pe.space_before = Pt(0)
    re_ = pe.add_run()
    re_.text = ("•  " if bullet else "") + en
    _style_run(re_, en_size, en_color, bold=bold)
    pj = tf.add_paragraph()
    pj.alignment = PP_ALIGN.LEFT
    pj.space_after = Pt(space_after)
    pj.space_before = Pt(0)
    rj = pj.add_run()
    rj.text = ("　 " if bullet else "") + jp
    _style_run(rj, jp_size, GREY, font=JP_FONT)
    return pe, pj


def card(slide, x, y, w, h, bg=CARD_BG, border=BORDER):
    """Flat professional box: light fill, thin full border, even padding. No stripes."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = border
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.11)
    tf.margin_bottom = Inches(0.11)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    return sh, tf


def label(slide, text):
    """Single section label, top-left, navy small-caps, letter-spaced."""
    tb, tf = textbox(slide, MX, Inches(0.4), CW, Inches(0.34))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _style_run(r, 13, NAVY, bold=True)
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", "200")
    return tb


def title(slide, en, jp, size=32, color=NAVY, jp_size=15, y=0.84):
    tb, tf = textbox(slide, MX, Inches(y), CW, Inches(1.15))
    pe = tf.paragraphs[0]
    pe.space_after = Pt(2)
    re_ = pe.add_run()
    re_.text = en
    _style_run(re_, size, color, bold=True, font=HEAD_FONT)
    pj = tf.add_paragraph()
    pj.space_before = Pt(0)
    rj = pj.add_run()
    rj.text = jp
    _style_run(rj, jp_size, GREY, font=JP_FONT)
    return tb


def infobox(slide, x, y, w, h, en, jp, bg=TINT):
    """Light box for an honesty / guard note. No edge stripe, no warning red."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = BORDER
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.09)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    en_jp(tf, en, jp, en_size=14, jp_size=11.5, en_color=NAVY, first=True, space_after=0)
    return sh


def stat_callout(slide, x, y, w, value, en_label, jp_label, vsize=42):
    _, tf = card(slide, x, y, w, Inches(1.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pv = tf.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    pv.space_after = Pt(2)
    rv = pv.add_run()
    rv.text = value
    _style_run(rv, vsize, NAVY, bold=True, font=HEAD_FONT)
    ple = tf.add_paragraph()
    ple.alignment = PP_ALIGN.CENTER
    ple.space_after = Pt(0)
    rle = ple.add_run()
    rle.text = en_label
    _style_run(rle, 12.5, INK, bold=True)
    plj = tf.add_paragraph()
    plj.alignment = PP_ALIGN.CENTER
    plj.space_before = Pt(0)
    rlj = plj.add_run()
    rlj.text = jp_label
    _style_run(rlj, 11, GREY, font=JP_FONT)


def caption(slide, x, y, w, en, jp):
    tb, tf = textbox(slide, x, y, w, Inches(0.75))
    pe = tf.paragraphs[0]
    pe.space_after = Pt(0)
    re_ = pe.add_run()
    re_.text = en
    _style_run(re_, 12, GREY, italic=True)
    pj = tf.add_paragraph()
    pj.space_before = Pt(0)
    rj = pj.add_run()
    rj.text = jp
    _style_run(rj, 10.5, GREY, italic=True, font=JP_FONT)


def placeholder(slide, x, y, w, h, en="image", jp="画像"):
    """Empty dashed-border box for the user to drop a photo / logo into later."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    sh.line.color.rgb = RGBColor(0xB2, 0xBB, 0xC8)
    sh.line.width = Pt(1.0)
    sh.shadow.inherit = False
    ln = sh.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"[ {en} ]"
    _style_run(r, 12.5, RGBColor(0x97, 0xA1, 0xAE))
    pj = tf.add_paragraph()
    pj.alignment = PP_ALIGN.CENTER
    rj = pj.add_run()
    rj.text = jp
    _style_run(rj, 10.5, RGBColor(0x97, 0xA1, 0xAE), font=JP_FONT)
    return sh


def picture_fit(slide, key, x, y, max_w, max_h):
    from PIL import Image
    png = PNG[key]
    iw, ih = Image.open(png).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    slide.shapes.add_picture(str(png), Emu(int(px)), Emu(int(py)), Emu(int(w)), Emu(int(h)))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def foot(slide, num):
    """Page number only, bottom-right, grey."""
    tb, tf = textbox(slide, Inches(11.0), Inches(7.04), Inches(1.83), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{num} / {N_SLIDES}"
    _style_run(r, 10, GREY)


def _cell(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None,
          jp=None, font=EN_FONT):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _style_run(r, size, color, bold=bold, font=font)
    if jp:
        pj = tf.add_paragraph()
        pj.alignment = align
        rj = pj.add_run()
        rj.text = jp
        _style_run(rj, size - 3.5, GREY, font=JP_FONT)


# =============================================================================
def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    def new():
        s = prs.slides.add_slide(blank)
        _set_bg(s, WHITE)
        return s

    # ---- 1 TITLE ----
    s = new()
    tb, tf = textbox(s, MX, Inches(0.9), CW, Inches(0.5))
    r = tf.paragraphs[0].add_run()
    r.text = "福井大学 ・ 竹本研 ・ 福井観光PBL"
    _style_run(r, 17, NAVY, bold=True, font=JP_FONT)
    tb, tf = textbox(s, MX, Inches(1.85), CW, Inches(1.7))
    pe = tf.paragraphs[0]
    pe.space_after = Pt(3)
    r = pe.add_run()
    r.text = "From Hokuriku reviews to tourism nudges"
    _style_run(r, 38, NAVY, bold=True, font=HEAD_FONT)
    pj = tf.add_paragraph()
    pj.space_before = Pt(2)
    rj = pj.add_run()
    rj.text = "北陸の口コミから、観光ナッジへ"
    _style_run(rj, 18, GREY, font=JP_FONT)
    tb2, tf2 = textbox(s, MX, Inches(3.95), CW, Inches(1.1))
    en_jp(tf2,
          "Which visitor problems can better information address? Which problems require operator action?",
          "より良い情報で改善できる問題は何か。事業者の対応が必要な問題は何か。",
          en_size=16, jp_size=13, en_color=INK, first=True, space_after=0)
    # presenters + date, formal, no chips
    tbp, tfp = textbox(s, MX, Inches(5.55), CW, Inches(0.9))
    pe = tfp.paragraphs[0]
    pe.space_after = Pt(2)
    r = pe.add_run()
    r.text = "Green ANDREW  ・  Xu ZILIN"
    _style_run(r, 18, NAVY, bold=True)
    pj = tfp.add_paragraph()
    rj = pj.add_run()
    rj.text = "2026年6月30日"
    _style_run(rj, 14, GREY, font=JP_FONT)
    # scope stat line (traced)
    sc_tb, sc_tf = textbox(s, MX, Inches(6.55), CW, Inches(0.8))
    en_jp(sc_tf,
          f"{TOTAL_REVIEWS} Google reviews across {N_POIS} sites in Hokuriku ●",
          f"北陸{N_POIS}スポットのGoogle口コミ{TOTAL_REVIEWS}件 ●",
          en_size=12, jp_size=10.5, en_color=GREY, first=True, space_after=0)
    notes(s, "Andrew (35 sec): Good morning. We studied online reviews of tourism sites in Hokuriku. Our question is "
              "simple: which visitor problems may be reduced through better information, and which problems require "
              "changes by the operator? We use the results to choose future experiments. We do not claim that any "
              "nudge already works.")

    # ---- 2 INTRODUCTION ----
    s = new()
    label(s, "I.  INTRODUCTION")
    title(s, "Research question", "研究の問い")
    _, tf = textbox(s, MX, Inches(2.15), Inches(7.5), Inches(4.0))
    en_jp(tf, "Reviews describe visitor pain points and reasons to visit.",
          "口コミには、訪問者の不満点と訪問したい理由が表れる。",
          first=True, bullet=True, space_after=18)
    en_jp(tf, "Some pain points may be reduced with clearer information before a trip.",
          "一部の不満点は、旅行前の分かりやすい情報で減らせる可能性がある。",
          bullet=True, space_after=18)
    en_jp(tf, "We rank low-cost ideas that should be tested next.",
          "次に検証すべき低コストの施策を順位づける。",
          bullet=True, space_after=0)
    levers = [("Information provision", "情報提供"),
              ("Pre-commitment", "事前コミットメント"),
              ("Demand redistribution", "需要の再配分")]
    cy = Inches(2.05)
    for i, (en, jp) in enumerate(levers):
        _, ctf = card(s, Inches(8.6), cy + Inches(i * 1.3), Inches(4.0), Inches(1.12))
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pe = ctf.paragraphs[0]
        pe.alignment = PP_ALIGN.CENTER
        pe.space_after = Pt(1)
        rr = pe.add_run()
        rr.text = en
        _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
        pjp = ctf.add_paragraph()
        pjp.alignment = PP_ALIGN.CENTER
        rjp = pjp.add_run()
        rjp.text = jp
        _style_run(rjp, 12, GREY, font=JP_FONT)
    infobox(s, MX, Inches(6.3), Inches(7.5), Inches(0.7),
            "Exploratory study: the results identify experiments to test. They do not prove cause and effect.",
            "探索的研究:結果は検証候補を示すものであり、因果関係を証明するものではない。")
    foot(s, 2)
    notes(s, "Andrew (40 sec): Reviews tell us about pain points, such as unclear opening hours, and attractions, such "
              "as scenery. Some problems may respond to information before the visit. Other problems, such as price or "
              "cleanliness, need action by the operator. This study ranks practical ideas for future testing.")

    # ---- 3 METHODS DATA ----
    s = new()
    label(s, "II.  METHODS")
    title(s, "The data", "データ")
    _, c1 = card(s, MX, Inches(2.05), Inches(5.7), Inches(2.05))
    pe = c1.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Tagged reviews"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = c1.add_paragraph()
    rj = pj.add_run()
    rj.text = "タグ付き口コミ"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c1, f"{TAGGED_ROWS} Google reviews", f"Google口コミ {TAGGED_ROWS}件", en_size=17, jp_size=12, en_color=NAVY, bold=True)
    en_jp(c1, f"Japanese {LANG_JP} · English {LANG_EN} · Chinese {LANG_CN}",
          f"日本語 {LANG_JP} ・ 英語 {LANG_EN} ・ 中国語 {LANG_CN}",
          en_size=14.5, jp_size=11.5, en_color=INK, space_after=0)
    _, c2 = card(s, MX, Inches(4.3), Inches(5.7), Inches(1.95))
    pe = c2.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Sites by prefecture"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = c2.add_paragraph()
    rj = pj.add_run()
    rj.text = "県別スポット数"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c2, f"Fukui {POI_FUKUI} · Ishikawa {POI_ISHIKAWA} · Toyama {POI_TOYAMA}",
          f"福井 {POI_FUKUI} ・ 石川 {POI_ISHIKAWA} ・ 富山 {POI_TOYAMA}",
          en_size=15, jp_size=12)
    en_jp(c2, f"Other languages {LANG_OTHER} · undetected or too short {LANG_UNDETECTED}",
          f"その他の言語 {LANG_OTHER} ・ 判定不能または短文 {LANG_UNDETECTED}",
          en_size=13, jp_size=10.5, en_color=GREY, space_after=0)
    picture_fit(s, "volume", Inches(6.55), Inches(2.0), Inches(6.2), Inches(3.5))
    caption(s, Inches(6.65), Inches(5.6), Inches(6.0),
            "The rating model uses Japanese, English, and Chinese-language Google reviews.",
            "評価モデルでは、日本語・英語・中国語のGoogle口コミを使用。")
    foot(s, 3)
    notes(s, f"Lynn (45 sec): Our dataset contains {TAGGED_ROWS} Google reviews from {N_POIS} sites in Fukui, "
              f"Ishikawa, and Toyama. The rating model uses {LANG_JP} Japanese-language reviews, {LANG_EN} English-language "
              f"reviews, and {LANG_CN} Chinese-language reviews. Language means the language of the text. It does not "
              f"tell us the reviewer's nationality. Later, I will show a separate set of {CN_ROWS} Xiaohongshu posts.")

    # ---- 4 METHODS PIPELINE ----
    s = new()
    label(s, "II.  METHODS")
    title(s, "How review text became data", "口コミテキストをデータに変える方法")
    _, tf = textbox(s, MX, Inches(2.15), Inches(7.6), Inches(4.0))
    en_jp(tf, f"1. Tag each review with {N_ASPECTS} aspects using human-reviewed keyword lists.",
          f"1. 人が確認したキーワードリストを使い、各口コミに{N_ASPECTS}の項目を付ける。",
          first=True, bullet=True, space_after=17)
    en_jp(tf, f"2. Define a low rating as {LOW_RATING_STARS} stars or fewer: {LOW_RATING_ROWS} of {MODEL_ROWS} reviews.",
          f"2. {LOW_RATING_STARS}つ星以下を低評価と定義:{MODEL_ROWS}件中{LOW_RATING_ROWS}件。",
          bullet=True, space_after=17)
    en_jp(tf, "3. Test whether each pain-point tag is associated with a low rating.",
          "3. 各不満点への言及が低評価と関連するかを検証する。",
          bullet=True, space_after=17)
    en_jp(tf, "Google star ratings provide the shared outcome. Sentiment-tool scores are not compared across languages.",
          "共通の結果変数にはGoogleの星評価を使う。感情分析ツールの点数は言語間で比較しない。",
          bullet=True, space_after=0)
    stat_callout(s, Inches(8.55), Inches(2.2), Inches(4.05), MODEL_ROWS, "modelled rows", "モデル対象行")
    stat_callout(s, Inches(8.55), Inches(4.05), Inches(4.05), LOW_RATING_ROWS, "low-rated rows", "低評価の行")
    foot(s, 4)
    notes(s, f"Andrew (50 sec): We used human-reviewed keyword lists to tag {MODEL_ROWS} Japanese, English, and "
              f"Chinese-language Google reviews. A low rating means {LOW_RATING_STARS} stars or fewer. There were {LOW_RATING_ROWS} "
              "low-rated reviews. For each pain point, we tested whether a review mentioning it was more likely to have "
              "a low rating. We use star ratings as the common outcome, so different sentiment tools are not compared.")

    # ---- 5 METHODS MODEL (stats-heavy) ----
    s = new()
    label(s, "II.  METHODS")
    title(s, "The statistical model", "統計モデル")
    _, tf = textbox(s, MX, Inches(2.05), Inches(7.0), Inches(4.0))
    en_jp(tf, "Model: Firth bias-reduced logistic regression, fitted separately for each aspect.",
          "モデル:各項目についてFirthバイアス低減ロジスティック回帰を実施。",
          first=True, bullet=True, en_size=16, jp_size=12.5, space_after=15)
    en_jp(tf, "Why Firth: several pain points were mentioned rarely. Firth is more stable with sparse data.",
          "Firthを使う理由:言及数が少ない不満点でも、より安定して推定できる。",
          bullet=True, en_size=16, jp_size=12.5, space_after=15)
    en_jp(tf, "Adjusted for text length, review language, and prefecture.",
          "口コミの長さ、言語、県を調整。",
          bullet=True, en_size=16, jp_size=12.5, space_after=15)
    en_jp(tf, f"Output: adjusted odds ratio, 95% confidence interval, and p-value for each of {MODELS_PRIMARY} primary models.",
          f"出力:{MODELS_PRIMARY}の主モデルごとに、調整済みオッズ比・95%信頼区間・p値を算出。",
          bullet=True, en_size=16, jp_size=12.5, space_after=0)
    model_steps = [
        ("Outcome", "低評価", f"{LOW_RATING_STARS} stars or fewer", f"{LOW_RATING_STARS}つ星以下"),
        ("Predictor", "説明変数", "One aspect tag", "1つの項目タグ"),
        ("Controls", "調整項目", "Length, language, prefecture", "長さ・言語・県"),
    ]
    for i, (en_h, jp_h, en_b, jp_b) in enumerate(model_steps):
        _, ctf = card(s, Inches(8.2), Inches(2.05 + i * 1.38), Inches(4.45), Inches(1.12))
        en_jp(ctf, en_h, jp_h, en_size=15, jp_size=11, en_color=NAVY, bold=True, first=True, space_after=2)
        en_jp(ctf, en_b, jp_b, en_size=13, jp_size=10.5, en_color=INK, space_after=0)
    foot(s, 5)
    notes(s, "Andrew (50 sec): We used Firth logistic regression because several pain points were mentioned only a "
              "small number of times. This method reduces small-sample bias and is more stable than ordinary logistic "
              "regression when data are sparse. Each model adjusts for review length, language, and prefecture. The "
              "result is an adjusted odds ratio with a confidence interval.")

    # ---- 6 METHODS INFERENCE + DECISION RULE (stats-heavy) ----
    s = new()
    label(s, "II.  METHODS")
    title(s, "Rules used to avoid overclaiming", "過大な主張を避けるためのルール")
    _, tf = textbox(s, MX, Inches(2.05), Inches(6.95), Inches(4.4))
    en_jp(tf, "Correct p-values for multiple tests using Benjamini-Hochberg FDR.",
          "複数の検定にはBenjamini-Hochberg FDR補正を行う。",
          first=True, bullet=True, en_size=16, jp_size=12.5, space_after=16)
    en_jp(tf, f"Ranking requires a harmful association, FDR significance, and at least {MIN_RANKING_MENTIONS} pooled mentions.",
          f"順位づけには、有害な関連、FDR有意、統合データで{MIN_RANKING_MENTIONS}件以上の言及が必要。",
          bullet=True, en_size=16, jp_size=12.5, space_after=16)
    en_jp(tf, f"For {FIRTH_SANITY_ASPECTS} selected aspects, Firth and standard-logit estimates were similar (maximum |Δlog OR| = {FIRTH_SANITY}).",
          f"選択した{FIRTH_SANITY_ASPECTS}項目では、Firthと通常ロジットの推定値は近かった(最大|Δlog OR|={FIRTH_SANITY})。",
          bullet=True, en_size=16, jp_size=12.5, space_after=16)
    en_jp(tf, "Classify each result as an information nudge or an operator fix.",
          "各結果を、情報ナッジまたは事業者による改善に分類する。",
          bullet=True, en_size=16, jp_size=12.5, space_after=0)
    _, guard = card(s, Inches(8.05), Inches(2.05), Inches(4.55), Inches(4.25), bg=TINT)
    guard.vertical_anchor = MSO_ANCHOR.MIDDLE
    en_jp(guard, "Statistical evidence", "統計的根拠", en_size=18, jp_size=13, en_color=NAVY, bold=True, first=True, space_after=10)
    en_jp(guard, "Association + uncertainty + enough mentions", "関連性・不確実性・十分な言及数", en_size=15, jp_size=11.5, space_after=18)
    en_jp(guard, "Action check", "実行可能性の確認", en_size=18, jp_size=13, en_color=NAVY, bold=True, space_after=10)
    en_jp(guard, "Can information help, or must the site change?", "情報で改善できるか、現地の変更が必要か。", en_size=15, jp_size=11.5, space_after=0)
    foot(s, 6)
    notes(s, f"Andrew (50 sec): We applied four safeguards. First, BH-FDR correction reduces false positives from "
              f"multiple tests. Second, a ranked opportunity needs at least {MIN_RANKING_MENTIONS} pooled mentions. "
              "Third, the association must point toward low ratings. Fourth, we separate problems that information may "
              "address from problems that require an operator change. These rules reduce cherry-picking.")

    # ---- 7 RESULTS: stat table ----
    s = new()
    label(s, "III.  RESULTS")
    title(s, "Pain points associated with low ratings", "低評価と関連する不満点")
    _, tf = textbox(s, MX, Inches(1.9), CW, Inches(0.55))
    en_jp(tf, "Two clear information priorities and one preliminary signal.",
          "明確な情報改善の優先項目は2つ。予備的なシグナルは1つ。",
          first=True, en_size=15, jp_size=11.5, space_after=0)
    headers = ["Pain point", "Mentions", "OR (95% CI)", "FDR p", "Interpretation"]
    body = [
        ("Opening hours", "営業時間", OPEN_N, OPEN_OR, f"{OPEN_CIL} to {OPEN_CIH}", OPEN_PV, "Priority", True),
        ("Itinerary / time", "旅程・所要時間", TIME_N, TIME_OR, f"{TIME_CIL} to {TIME_CIH}", TIME_PV, "Priority", True),
        ("Wayfinding", "道案内", SIGN_N, SIGN_OR, f"{SIGN_CIL} to {SIGN_CIH}", SIGN_PV, "Preliminary", False),
    ]
    n_rows = len(body) + 1
    tx, ty, tw, th = MX, Inches(2.7), Inches(8.25), Inches(2.35)
    gt = s.shapes.add_table(n_rows, 5, tx, ty, tw, th).table
    gt.first_row = False
    gt.horz_banding = False
    widths = [Inches(2.35), Inches(1.0), Inches(2.15), Inches(1.05), Inches(1.7)]
    for i, wv in enumerate(widths):
        gt.columns[i].width = wv
    for j, htext in enumerate(headers):
        _cell(gt.cell(0, j), htext, 12.5, WHITE, bold=True,
              align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=NAVY)
    for ri, (en, jp, mentions, orv, ci, pv, interpretation, ok) in enumerate(body, start=1):
        rowfill = WHITE if ri % 2 else CARD_BG
        _cell(gt.cell(ri, 0), en, 12, INK, bold=True, fill=rowfill, jp=jp)
        _cell(gt.cell(ri, 1), mentions, 12, INK, align=PP_ALIGN.CENTER, fill=rowfill)
        _cell(gt.cell(ri, 2), f"{orv}  ({ci})", 12, INK, align=PP_ALIGN.CENTER, fill=rowfill)
        _cell(gt.cell(ri, 3), pv, 12, INK, align=PP_ALIGN.CENTER, fill=rowfill)
        _cell(gt.cell(ri, 4), interpretation, 12, NAVY if ok else GREY, bold=ok, align=PP_ALIGN.CENTER, fill=rowfill)
    picture_fit(s, "nudge_aspect_fig", Inches(9.0), Inches(2.55), Inches(3.75), Inches(3.2))
    infobox(s, MX, Inches(5.35), Inches(8.25), Inches(1.0),
            f"Price (OR {PRICE_OR}) and cleanliness (OR {CLEAN_OR}) were also linked to low ratings, but they require operator action.",
            f"価格(OR {PRICE_OR})と清潔さ(OR {CLEAN_OR})も低評価と関連したが、事業者の対応が必要。")
    caption(s, MX, Inches(6.45), Inches(8.25),
            "OR = adjusted odds ratio. Wayfinding is borderline: its 95% interval includes 1.",
            "OR=調整済みオッズ比。道案内は境界的で、95%信頼区間に1を含む。")
    foot(s, 7)
    notes(s, f"Andrew (60 sec): Opening-hour problems were mentioned in {OPEN_N} reviews and had an adjusted odds ratio "
              f"of {OPEN_OR}. Itinerary and time problems were mentioned in only {TIME_N} reviews, but their association "
              f"was also strong, with an odds ratio of {TIME_OR}. Wayfinding had an odds ratio of {SIGN_OR}, but its "
              "confidence interval includes one. We therefore call wayfinding preliminary. These are associations, not "
              "causal effects.")

    # ---- 8 RESULTS: the two justified nudges ----
    s = new()
    label(s, "III.  RESULTS")
    title(s, "Two priorities and one idea to test", "優先すべき2項目と、検証する1案")
    _, b1 = card(s, MX, Inches(2.0), Inches(7.55), Inches(2.0))
    pe = b1.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Opening hours and availability"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = b1.add_paragraph()
    rj = pj.add_run()
    rj.text = "開館時間・営業状況"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(b1, f"Reviews mentioning opening-hour problems were more likely to give {LOW_RATING_STARS} stars or fewer. Show current hours, last entry, and possible closures before visitors travel.",
          f"営業時間の問題に触れた口コミでは、{LOW_RATING_STARS}つ星以下の評価が多く見られた。出発前に、営業時間・最終入場時刻・臨時休業の可能性を知らせる。",
          en_size=15, jp_size=11.5, space_after=0)
    _, b2 = card(s, MX, Inches(4.15), Inches(7.55), Inches(2.05))
    pe = b2.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Itinerary fit and time-cost"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = b2.add_paragraph()
    rj = pj.add_run()
    rj.text = "行程適合・所要時間"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(b2, f"Reviews mentioning timing or itinerary problems were more likely to give {LOW_RATING_STARS} stars or fewer. Show realistic visit times and route order before plans are finalized.",
          f"時間や旅程の問題に触れた口コミでは、{LOW_RATING_STARS}つ星以下の評価が多く見られた。計画を決める前に、現実的な滞在時間と観光ルートを示す。",
          en_size=15, jp_size=11.5, space_after=0)
    picture_fit(s, "nudge_info_fig", Inches(8.35), Inches(2.0), Inches(4.4), Inches(4.2))
    infobox(s, MX, Inches(6.35), Inches(12.13), Inches(0.65),
            f"Idea to test: clearer route guidance. Evidence is limited (n = {SIGN_N}; OR {SIGN_OR}, 95% CI {SIGN_CIL} to {SIGN_CIH}).",
            f"検証する案:より分かりやすいルート案内。データは限定的(n={SIGN_N}; OR {SIGN_OR}, 95%CI {SIGN_CIL}〜{SIGN_CIH})。")
    foot(s, 8)
    notes(s, "Andrew (55 sec): Our first priority is a simple visit-readiness card showing opening hours, last entry, "
              "and closure risks. Our second priority is realistic visit duration and route order before the visitor "
              "finalizes a plan. Wayfinding remains an idea to test because the evidence is limited. These proposals "
              "follow the observed associations, but only an experiment can show whether they work.")

    # ---- 9 RESULTS: Chinese-language Xiaohongshu context ----
    s = new()
    label(s, "III.  RESULTS")
    title(s, "What Chinese-language Xiaohongshu posts discuss", "中国語の小紅書投稿で話題になったこと")
    _, c1 = card(s, MX, Inches(2.0), Inches(5.7), Inches(3.15))
    pe = c1.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Separate, directional evidence"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = c1.add_paragraph()
    rj = pj.add_run()
    rj.text = "別枠の方向性を示すデータ"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c1, f"{CN_ROWS} Chinese-language posts about Fukui, all from Xiaohongshu.",
          f"福井に関する中国語投稿{CN_ROWS}件。すべて小紅書から収集。",
          en_size=14, jp_size=11, space_after=12)
    en_jp(c1, "These posts have no star ratings and do not enter the Google rating model.",
          "星評価がないため、Google評価モデルには含めない。",
          en_size=13.5, jp_size=10.8, en_color=GREY, space_after=0)
    _, c2 = card(s, Inches(6.55), Inches(2.0), Inches(6.1), Inches(3.15))
    pe = c2.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Topic mentions"
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = c2.add_paragraph()
    rj = pj.add_run()
    rj.text = "話題への言及"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c2, f"Scenic nature appeared in {XHS_SCENIC_N} of {CN_ROWS} posts ({XHS_SCENIC_PCT}).",
          f"自然景観は{CN_ROWS}件中{XHS_SCENIC_N}件({XHS_SCENIC_PCT})で言及。",
          en_size=14, jp_size=11, space_after=12)
    en_jp(c2, f"Dinosaurs and museums appeared in {XHS_DINO_N} posts ({XHS_DINO_PCT}).",
          f"恐竜・博物館は{XHS_DINO_N}件({XHS_DINO_PCT})で言及。",
          en_size=14, jp_size=11, space_after=0)
    infobox(s, MX, Inches(5.45), Inches(12.05), Inches(0.92),
            "Test idea: compare Chinese-language discovery cards featuring scenic nature and dinosaurs/museums. Measure clicks, saves, and itinerary additions.",
            "検証案:自然景観と恐竜・博物館を紹介する中国語カードを比較し、クリック・保存・旅程への追加を測定する。")
    _, tf = textbox(s, MX, Inches(6.48), Inches(12.0), Inches(0.42))
    en_jp(tf, "Topic presence only: a mention does not show positive or negative opinion, importance, or causal impact.",
          "話題の有無のみ:言及は、肯定・否定、重要度、因果効果を示すものではない。",
          first=True, en_size=11.5, jp_size=9.5, en_color=GREY, space_after=0)
    foot(s, 9)
    notes(s, f"Lynn (60 sec): We also analyzed {CN_ROWS} Chinese-language Xiaohongshu posts about Fukui. These posts "
              f"have no star ratings, so we keep them separate from the Google model. Scenic nature appeared in "
              f"{XHS_SCENIC_N} posts, and dinosaurs or museums appeared in {XHS_DINO_N}. These are topic mentions only. "
              "They suggest content for a Chinese-language discovery-card test, but they do not show sentiment or impact.")

    # ---- 10 RESULTS: POI action map ----
    s = new()
    label(s, "III.  RESULTS")
    title(s, "Candidate sites for action", "施策候補となるスポット")
    stat_callout(s, MX, Inches(1.95), Inches(2.05), FIX_COUNT, "Fix-it sites", "改善型", vsize=36)
    stat_callout(s, Inches(2.78), Inches(1.95), Inches(2.05), PROMOTE_COUNT, "Promote-it sites", "推奨型", vsize=36)
    stat_callout(s, Inches(4.96), Inches(1.95), Inches(2.05), CROWD_COUNT, "Crowding hotspots", "混雑ホット", vsize=36)
    _, tf = textbox(s, MX, Inches(3.85), Inches(6.5), Inches(3.0))
    en_jp(tf, "Fix-it sites have high review volume and unusually frequent pain-point tags.",
          "改善型は口コミ数が多く、不満点タグが通常より多いスポット。",
          first=True, en_size=14.5, jp_size=11.5, space_after=14)
    en_jp(tf, "Promote-it sites have strong star ratings and lower review volume.",
          "推奨型は星評価が高く、口コミ数が少ないスポット。",
          en_size=14.5, jp_size=11.5, en_color=GREY, space_after=14)
    en_jp(tf, f"All {PROMOTE_COUNT} promote-it sites are exploratory; {PROMOTE_STRICT_COUNT} met the strict rule. Review volume does not equal visitor volume.",
          f"推奨型{PROMOTE_COUNT}件はすべて探索的で、厳格基準を満たしたものは{PROMOTE_STRICT_COUNT}件。口コミ数は来訪者数ではない。",
          en_size=14, jp_size=11, space_after=0)
    picture_fit(s, "nudge_poi_fig", Inches(7.15), Inches(1.95), Inches(5.5), Inches(4.9))
    foot(s, 10)
    notes(s, f"Andrew (55 sec): We also created an exploratory site-level index. It identifies {FIX_COUNT} fix-it "
              f"candidates, {PROMOTE_COUNT} promote-it candidates, and {CROWD_COUNT} crowding hotspots. These are places "
              "for follow-up work, not final classifications. Promote-it uses high positive Google star-rating shares "
              f"ratings among lower-volume sites. No site met the strict promote-it rule. Fukui examples are {PROMO1}, "
              f"n={PROMO1_N}, {PROMO1_SHARE} positive, and {PROMO2}, n={PROMO2_N}, {PROMO2_SHARE} positive.")

    # ---- 11 RESULTS: FINAL CROSS-LANGUAGE PRIORITIES ----
    s = new()
    label(s, "III.  RESULTS")
    title(s, "Three candidate experiments, ranked",
          "3つの実験候補を順位づけ")
    _, tf = textbox(s, MX, Inches(1.82), CW, Inches(0.55))
    en_jp(tf,
          "These are three pre-specified candidates. Evidence strength differs by source and remains separate.",
          "事前に設定した3つの候補。根拠の強さはソースごとに異なり、別々に扱う。",
          first=True, en_size=13.5, jp_size=10.5, en_color=GREY, space_after=0)
    gt = s.shapes.add_table(4, 4, MX, Inches(2.48), Inches(12.05), Inches(3.92)).table
    widths = [Inches(0.75), Inches(2.55), Inches(4.35), Inches(4.4)]
    for i, width in enumerate(widths):
        gt.columns[i].width = width
    headers = [
        ("Rank", "順位"),
        ("Common solution", "共通施策"),
        ("Why this rank", "順位の理由"),
        ("Next-semester test", "来学期の実験"),
    ]
    for j, (en, jp) in enumerate(headers):
        _cell(gt.cell(0, j), en, 11.5, WHITE, bold=True, align=PP_ALIGN.CENTER, fill=NAVY, jp=jp)
    evidence_en = [
        "Strong Google rating associations for hours and itinerary; easiest prototype.",
        "Cross-language attraction signals; Chinese XHS topic evidence is directional.",
        "Crowding appears across sources; rating association remains preliminary; harder to implement.",
    ]
    evidence_ja = [
        "営業時間と旅程はGoogle評価と強く関連。試作が最も容易。",
        "言語をまたぐ魅力シグナル。中国語XHSの話題データは方向性のみ。",
        "複数ソースで混雑に言及。評価との関連は予備的で、実装も難しい。",
    ]
    for row_index, priority in enumerate(PRIORITIES, start=1):
        row_fill = WHITE if row_index % 2 else CARD_BG
        _cell(gt.cell(row_index, 0), priority["rank"], 16, NAVY, bold=True,
              align=PP_ALIGN.CENTER, fill=row_fill)
        _cell(gt.cell(row_index, 1), priority["name_en"], 11.5, INK, bold=True,
              fill=row_fill, jp=priority["name_ja"])
        _cell(
            gt.cell(row_index, 2),
            f"{priority['ease']}\n{evidence_en[row_index - 1]}",
            10.5,
            INK,
            fill=row_fill,
            jp=f"実装容易性 {priority['ease']}\n{evidence_ja[row_index - 1]}",
        )
        _cell(gt.cell(row_index, 3), priority["test_en"], 10.5, INK,
              fill=row_fill, jp=priority["test_ja"])
    infobox(
        s, MX, Inches(6.55), Inches(12.05), Inches(0.52),
        "Opportunity ranking only. It does not estimate how well any intervention will work.",
        "機会の順位づけのみ。各施策の効果を推定したものではない。",
    )
    foot(s, 11)
    notes(s, "Andrew (55 sec): We ranked three pre-specified experiments. The first is the multilingual visit-readiness "
              "card because opening hours and itinerary fit have the strongest relevant rating evidence, and the card is "
              "easy to prototype. The second tests localized discovery content. The third tests off-peak and alternative "
              "site prompts. Evidence from different sources is kept separate. This ranking does not estimate effectiveness.")

    # ---- 12 DISCUSSION ----
    s = new()
    label(s, "IV.  DISCUSSION")
    title(s, "Study limits", "研究の限界")
    cavs = [
        ("Observational review data cannot establish cause and effect.", "観察データである口コミから、因果関係は証明できない。"),
        ("Opportunity scores rank follow-up experiments. They do not measure intervention effectiveness.", "機会スコアは次の実験を順位づけるもので、施策の効果を測るものではない。"),
        ("The Firth model does not account for reviews nested within each site. Uncertainty may be understated.", "Firthモデルはスポット内の口コミのまとまりを考慮していない。不確実性を過小評価する可能性がある。"),
        ("Language groups describe the review text, not reviewer nationality.", "言語グループは口コミの言語を示し、投稿者の国籍を示さない。"),
    ]
    _, tf = textbox(s, MX, Inches(2.05), Inches(8.1), Inches(4.2))
    for i, (en, jp) in enumerate(cavs):
        en_jp(tf, en, jp, first=(i == 0), bullet=True, en_size=15.5, jp_size=12, space_after=15)
    _, ltf = card(s, Inches(9.05), Inches(2.05), Inches(3.55), Inches(3.55), bg=TINT)
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    en_jp(ltf, "Main interpretation", "主な解釈", en_size=18, jp_size=13, en_color=NAVY, bold=True, first=True, space_after=12)
    en_jp(ltf, "Useful for choosing experiments", "実験候補の選択に役立つ", en_size=16, jp_size=12, space_after=18)
    en_jp(ltf, "Requires randomized testing before policy use", "政策に使う前に無作為化実験が必要", en_size=16, jp_size=12, space_after=0)
    foot(s, 12)
    notes(s, "Andrew (45 sec): These results have important limits. Reviews are observational and self-selected, so "
              "the associations are not causal. The model also treats review rows separately and does not model clustering "
              "within tourism sites. This may understate uncertainty. Language labels describe the text, not nationality. "
              "The results are useful for choosing experiments, but randomized testing is required before policy use.")

    # ---- 13 FUTURE WORK (close, white) ----
    s = new()
    label(s, "IV.  DISCUSSION")
    title(s, "Next step: test the first nudge", "次のステップ:第1候補を検証")
    _, c1 = card(s, MX, Inches(2.1), Inches(6.0), Inches(2.55))
    pe = c1.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Pre-register the test"
    _style_run(rr, 17, NAVY, bold=True, font=HEAD_FONT)
    pj = c1.add_paragraph()
    rj = pj.add_run()
    rj.text = "実験計画を事前登録"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c1, "Record the hypothesis, target site, randomization unit, and outcomes before data collection.",
          "データ収集前に、仮説・対象スポット・無作為化単位・評価指標を記録する。",
          en_size=13.5, jp_size=11, space_after=4)
    pe = c1.add_paragraph()
    rr = pe.add_run()
    rr.text = "docs/nudge_experiment_register.html"
    _style_run(rr, 11, BLUE, font="Consolas")
    _, c2 = card(s, Inches(6.95), Inches(2.1), Inches(5.7), Inches(2.55))
    pe = c2.paragraphs[0]
    rr = pe.add_run()
    rr.text = "First experiment"
    _style_run(rr, 17, NAVY, bold=True, font=HEAD_FONT)
    pj = c2.add_paragraph()
    rj = pj.add_run()
    rj.text = "最初の実験"
    _style_run(rj, 11.5, GREY, font=JP_FONT)
    en_jp(c2, f"Start with priority {PRIORITIES[0]['rank']}: {PRIORITIES[0]['name_en']}. Randomize exposure and measure clicks, saves, and itinerary additions.",
          f"優先順位{PRIORITIES[0]['rank']}の「{PRIORITIES[0]['name_ja']}」から開始。表示を無作為化し、クリック・保存・旅程への追加を測定する。",
          en_size=13.5, jp_size=11, space_after=0)
    infobox(
        s, MX, Inches(4.78), Inches(12.05), Inches(0.52),
        f"In one line: {TOTAL_REVIEWS} reviews across {N_POIS} POIs -> {FIX_COUNT} fix-it and "
        f"{PROMOTE_COUNT} promote-it sites -> 3 cross-language nudges.",
        f"一言で:{N_POIS}スポットの{TOTAL_REVIEWS}件の口コミ → {FIX_COUNT}の改善型・"
        f"{PROMOTE_COUNT}の推奨型 → 3つの言語横断ナッジ。",
    )
    cl_tb, cl_tf = textbox(s, MX, Inches(5.55), CW, Inches(1.3))
    pe = cl_tf.paragraphs[0]
    pe.space_after = Pt(3)
    rr = pe.add_run()
    rr.text = ("We mapped where Hokuriku visitors hit friction and which gems to promote. "
               "The analysis chooses the experiment; randomization tests whether it works.")
    _style_run(rr, 16, NAVY, bold=True, font=HEAD_FONT)
    pj = cl_tf.add_paragraph()
    rj = pj.add_run()
    rj.text = "北陸の旅行者がどこで不満を感じ、どの名所を推すべきかを地図化した。分析で実験を選び、無作為化で効果を検証する。"
    _style_run(rj, 12, GREY, font=JP_FONT)
    pe2 = cl_tf.add_paragraph()
    pe2.space_before = Pt(8)
    rr = pe2.add_run()
    rr.text = "Thank you.  ありがとうございました。"
    _style_run(rr, 15, NAVY, bold=True, font=JP_FONT)
    foot(s, 13)
    notes(s, f"Andrew (40 sec): Our next step is a pre-registered A/B test of priority {PRIORITIES[0]['rank']}, the "
              f"{PRIORITIES[0]['name_en']}. We will randomize exposure by visitor session and measure clicks, saves, "
              "and itinerary additions. The current analysis selects the experiment. Randomization will test whether "
              "the nudge changes behavior. Thank you.")

    prs.save(str(OUT_PPTX))
    return prs


def main() -> int:
    build()
    size = OUT_PPTX.stat().st_size
    print(f"wrote {OUT_PPTX} ({size:,} bytes); {N_SLIDES} slides; figures rasterized: {list(PNG)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
