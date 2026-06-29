#!/usr/bin/env python3
"""Build seminar speaker/Q&A notes and a two-page A4 academic handout.

All reported analysis numbers are resolved from aggregate CSV/JSON outputs.
The spoken script is read from the generated PowerPoint speaker notes.
"""
from __future__ import annotations

import json
import re
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "NUDGE-Seminar-Slides.pptx"
# The intro was split into two read-along slides (research question + approach)
# in NUDGE-Intro-Simple.pptx; its speaker notes replace the live deck's single
# "Research question" note for the Q&A script.
INTRO_PPTX = ROOT / "NUDGE-Intro-Simple.pptx"
# The presented deck is assembled from the simplified decks; methods slides 5-6
# collapse into one "How we tested it" slide, so the spoken script is sourced
# from these rather than the full live deck.
METHODS_PPTX = ROOT / "NUDGE-Methods-Simple.pptx"
RESULTS_PPTX = ROOT / "NUDGE-Results-Simple.pptx"
OUT_QA = ROOT / "Seminar-Speaker-Script-and-QA.docx"
OUT_HANDOUT = ROOT / "Seminar-Submission-Handout-A4.docx"

ASPECT_CSV = ROOT / "output/nudge_analysis/aspect_opportunity_map.csv"
ASPECT_MANIFEST = ROOT / "output/nudge_analysis/aspect_opportunity_map_manifest.json"
POI_CSV = ROOT / "output/nudge_analysis/poi_opportunity_index.csv"
POI_MANIFEST = ROOT / "output/nudge_analysis/poi_opportunity_index_manifest.json"
PRIORITY_CSV = ROOT / "output/nudge_analysis/cross_language_solution_priorities.csv"
PRIORITY_MANIFEST = ROOT / "output/nudge_analysis/cross_language_solution_priorities_manifest.json"
XHS_TOPICS_CSV = ROOT / "output/chinese_specific_insights_xhs_only/keyword_occurrence_by_category.csv"
XHS_TOPICS_MANIFEST = ROOT / "output/chinese_specific_insights_xhs_only/chinese_specific_insights_manifest.json"
CHINESE_GOOGLE_MANIFEST = ROOT / "output/chinese_google_reviews_analysis/tagged_chinese_google_reviews_manifest.json"

for path in (
    PPTX,
    ASPECT_CSV,
    ASPECT_MANIFEST,
    POI_CSV,
    POI_MANIFEST,
    PRIORITY_CSV,
    PRIORITY_MANIFEST,
    XHS_TOPICS_CSV,
    XHS_TOPICS_MANIFEST,
    CHINESE_GOOGLE_MANIFEST,
):
    if not path.exists():
        raise SystemExit(f"missing required input: {path}")

aspect_df = pd.read_csv(ASPECT_CSV)
aspect_manifest = json.loads(ASPECT_MANIFEST.read_text())
poi_df = pd.read_csv(POI_CSV)
poi_manifest = json.loads(POI_MANIFEST.read_text())
priority_df = pd.read_csv(PRIORITY_CSV)
xhs_df = pd.read_csv(XHS_TOPICS_CSV)
cn_google_manifest = json.loads(CHINESE_GOOGLE_MANIFEST.read_text())


def aspect(code: str) -> pd.Series:
    row = aspect_df[
        (aspect_df["analysis"] == "A_primary")
        & (aspect_df["segment"] == "pooled")
        & (aspect_df["aspect"] == code)
    ]
    if len(row) != 1:
        raise KeyError(f"expected one pooled primary row for {code}; got {len(row)}")
    return row.iloc[0]


def aspect_segment(code: str, segment: str) -> pd.Series:
    row = aspect_df[
        (aspect_df["analysis"] == "A_primary")
        & (aspect_df["segment"] == segment)
        & (aspect_df["aspect"] == code)
    ]
    if len(row) != 1:
        raise KeyError(f"expected one primary row for {code}/{segment}; got {len(row)}")
    return row.iloc[0]


def xhs_topic(code: str) -> pd.Series:
    row = xhs_df[(xhs_df["evidence_family"] == "topic") & (xhs_df["code"] == code)]
    if len(row) != 1:
        raise KeyError(f"expected one XHS topic row for {code}; got {len(row)}")
    return row.iloc[0]


def fmt_p(value: float) -> str:
    value = float(value)
    if value < 0.0001:
        return f"{value:.2e}".replace("e-0", "e-").replace("e+0", "e+")
    return f"{value:.4f}".rstrip("0").rstrip(".")


def pct(value: float, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%"


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_margins(cell, top=60, start=80, bottom=60, end=80) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def configure_a4(doc: Document, margin_cm: float = 1.7) -> None:
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(margin_cm)
    section.bottom_margin = Cm(margin_cm)
    section.left_margin = Cm(margin_cm)
    section.right_margin = Cm(margin_cm)


def set_default_styles(doc: Document, body_size: float = 10.5) -> None:
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(body_size)
    normal.font.color.rgb = RGBColor(0x14, 0x18, 0x1D)
    normal.paragraph_format.space_after = Pt(5)
    normal.paragraph_format.line_spacing = 1.07
    for name, size, color in (
        ("Title", 17, "16345E"),
        ("Heading 1", 14, "16345E"),
        ("Heading 2", 11.5, "16345E"),
    ):
        style = styles[name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)
        style.paragraph_format.space_before = Pt(7)
        style.paragraph_format.space_after = Pt(3)


def add_inline_section(doc: Document, label: str, text: str, size: float = 11.0) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.07
    lead = p.add_run(f"{label}. ")
    lead.bold = True
    lead.font.size = Pt(size)
    body = p.add_run(text)
    body.font.size = Pt(size)


def add_source_line(doc: Document, paths: list[Path]) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(1)
    p.paragraph_format.space_after = Pt(5)
    r = p.add_run("Sources: " + "; ".join(path.relative_to(ROOT).as_posix() for path in paths))
    r.font.name = "Arial"
    r.font.size = Pt(8)
    r.font.italic = True
    r.font.color.rgb = RGBColor(0x5A, 0x65, 0x73)


metrics = aspect_manifest["metrics"]
poi_metrics = poi_manifest["metrics"]
languages = metrics["tagged_language_group_counts"]
model_n = int(metrics["primary_model_rows"])
low_n = int(metrics["primary_low_rating_rows"])
total_n = int(metrics["tagged_input_rows"])
poi_n = int(poi_metrics["n_pois_total"])
min_mentions = int(aspect_manifest["filters"]["min_pooled_positive_for_ranking"])
n_aspects = int(metrics["A_primary_models_fit"])
sanity_aspects = int(len(metrics["plain_vs_firth_sanity"]))
low_rating_match = re.search(r"<=\s*(\d+)", aspect_manifest["filters"]["low_rating_definition"])
if not low_rating_match:
    raise ValueError("could not parse low-rating cutoff from analysis manifest")
low_rating_cutoff = int(low_rating_match.group(1))

opening = aspect("opening_hours_availability")
itinerary = aspect("itinerary_fit_time_cost")
wayfinding = aspect("wayfinding_signage")
price = aspect("price_value")
cleanliness = aspect("cleanliness_comfort")
waiting = aspect("waiting_crowding")
segment_counts = {
    code: {
        segment: int(aspect_segment(code, segment)["n_positive"])
        for segment in ("japanese", "english", "chinese")
    }
    for code in (
        "opening_hours_availability",
        "itinerary_fit_time_cost",
        "wayfinding_signage",
    )
}

scenic = xhs_topic("scenic_nature")
dinosaurs = xhs_topic("dinosaurs_museums")

promote_fukui = poi_df[
    poi_df["is_promote_it"].astype(bool) & poi_df["is_fukui"].astype(bool)
].sort_values("promote_it_score", ascending=False)
if len(promote_fukui) < 2:
    raise KeyError("expected at least two Fukui promote-it candidates")
promo1, promo2 = promote_fukui.iloc[0], promote_fukui.iloc[1]

priority1 = priority_df[priority_df["rank"] == 1]
if len(priority1) != 1:
    raise KeyError("expected one priority rank 1")
priority1 = priority1.iloc[0]

validation = cn_google_manifest["metrics"]["snownlp_validation"]


def build_handout() -> None:
    doc = Document()
    configure_a4(doc, 1.65)
    set_default_styles(doc, 11.0)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run("Turning Hokuriku Review Text into Testable Tourism Nudges")
    r.bold = True
    r.font.name = "Arial"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(0x16, 0x34, 0x5E)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(1)
    r = p.add_run("Green ANDREW and Xu ZILIN  |  Takemoto Lab, University of Fukui")
    r.bold = True
    r.font.size = Pt(10)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(7)
    r = p.add_run("June 30, 2026")
    r.font.size = Pt(9.5)

    add_inline_section(
        doc,
        "Research focus",
        "This exploratory study asks which tourism pain points may be reduced through clearer pre-visit information, "
        "which require operator action, and which low-cost ideas should be tested next. Review evidence is used to "
        "rank experiments. Causal effectiveness remains untested.",
    )
    add_inline_section(
        doc,
        "Data",
        f"The corpus contains {total_n:,} Google reviews from {poi_n} points of interest across Fukui, Ishikawa, and "
        f"Toyama. Language counts are {languages['japanese']:,} Japanese, {languages['english']:,} English, "
        f"{languages['chinese']:,} Chinese, {languages['other_non_english_non_japanese']:,} other-language, and "
        f"{languages['undetected_or_too_short']:,} undetected or too short. Language labels describe text language. "
        f"The primary model uses {model_n:,} Japanese, English, and Chinese-language reviews. Low rating means "
        f"{low_rating_cutoff} stars "
        f"or fewer; {low_n:,} modeled reviews meet that definition.",
    )
    add_inline_section(
        doc,
        "Method",
        f"Human-reviewed keyword codebooks tag {n_aspects} aspects. One Firth bias-reduced logistic regression is fitted per "
        f"aspect, with low rating as the outcome. Models adjust for text length, language, and prefecture. Firth was "
        f"chosen because several tags are rare. Benjamini-Hochberg false-discovery-rate correction addresses multiple "
        f"testing. Ranking also requires at least {min_mentions} pooled mentions. Odds ratios describe association in "
        f"odds and should not be read as probability ratios.",
    )

    heading = doc.add_paragraph()
    heading.paragraph_format.space_before = Pt(3)
    heading.paragraph_format.space_after = Pt(3)
    rr = heading.add_run("Main statistical findings")
    rr.bold = True
    rr.font.size = Pt(11)
    rr.font.color.rgb = RGBColor(0x16, 0x34, 0x5E)

    table = doc.add_table(rows=1, cols=5)
    table.autofit = False
    widths = (Cm(4.5), Cm(1.6), Cm(2.2), Cm(3.7), Cm(2.0))
    headers = ("Pain point", "Mentions", "Adjusted OR", "95% CI", "BH-FDR p")
    for i, (cell, text, width) in enumerate(zip(table.rows[0].cells, headers, widths)):
        cell.width = width
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        set_cell_shading(cell, "16345E")
        set_cell_margins(cell)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER if i else WD_ALIGN_PARAGRAPH.LEFT
        r = p.add_run(text)
        r.bold = True
        r.font.size = Pt(9.5)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    rows = (
        ("Opening hours", opening),
        ("Itinerary / time", itinerary),
        ("Wayfinding", wayfinding),
        ("Price / value", price),
        ("Cleanliness", cleanliness),
    )
    for idx, (label, row) in enumerate(rows):
        cells = table.add_row().cells
        values = (
            label,
            f"{int(row['n_positive']):,}",
            f"{float(row['odds_ratio']):.2f}",
            f"{float(row['or_ci_low']):.2f} to {float(row['or_ci_high']):.2f}",
            fmt_p(row["p_value_bh_fdr"]),
        )
        for i, (cell, text, width) in enumerate(zip(cells, values, widths)):
            cell.width = width
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            set_cell_margins(cell)
            if idx % 2:
                set_cell_shading(cell, "F2F5F9")
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER if i else WD_ALIGN_PARAGRAPH.LEFT
            r = p.add_run(text)
            r.font.size = Pt(9.5)
            if i == 0:
                r.bold = True

    add_inline_section(
        doc,
        "Interpretation",
        f"Opening hours (OR {opening['odds_ratio']:.2f}, 95% CI {opening['or_ci_low']:.2f} to "
        f"{opening['or_ci_high']:.2f}) and itinerary/time cost (OR {itinerary['odds_ratio']:.2f}, 95% CI "
        f"{itinerary['or_ci_low']:.2f} to {itinerary['or_ci_high']:.2f}) are the two clearest information priorities. "
        f"Wayfinding is preliminary: BH-FDR p={fmt_p(wayfinding['p_value_bh_fdr'])}, while its 95% CI "
        f"({wayfinding['or_ci_low']:.2f} to {wayfinding['or_ci_high']:.2f}) includes 1. Price and cleanliness show "
        f"associations with low ratings, but those issues require operator action.",
        size=10.8,
    )
    add_inline_section(
        doc,
        "Cross-language counts",
        f"Opening-hour mentions are Japanese {segment_counts['opening_hours_availability']['japanese']}, English "
        f"{segment_counts['opening_hours_availability']['english']}, and Chinese "
        f"{segment_counts['opening_hours_availability']['chinese']}. Itinerary/time mentions are "
        f"{segment_counts['itinerary_fit_time_cost']['japanese']}, "
        f"{segment_counts['itinerary_fit_time_cost']['english']}, and "
        f"{segment_counts['itinerary_fit_time_cost']['chinese']}. Wayfinding mentions are "
        f"{segment_counts['wayfinding_signage']['japanese']}, "
        f"{segment_counts['wayfinding_signage']['english']}, and "
        f"{segment_counts['wayfinding_signage']['chinese']}. Sparse counts make separate language-specific effect "
        f"estimates unreliable, so the primary analysis pools languages and adjusts for language.",
        size=10.8,
    )
    add_source_line(doc, [ASPECT_CSV, ASPECT_MANIFEST])

    doc.add_page_break()

    add_inline_section(
        doc,
        "Proposed nudges",
        f"Priority 1 is the {priority1['solution_label_en']}: show current hours, last entry, realistic visit duration, "
        f"transfers, booking needs, and backup options before plans are finalized. Priority 2 tests localized discovery "
        f"and itinerary cards. Priority 3 tests off-peak timing and alternative-site prompts. These are ordinal "
        f"opportunity rankings based on the configured evidence and implementation rules.",
    )
    add_inline_section(
        doc,
        "Chinese-language context",
        f"A separate dataset contains {int(scenic['denominator_posts']):,} Fukui-focused Xiaohongshu posts without star "
        f"ratings. Scenic nature appears in {int(scenic['count']):,} posts ({float(scenic['pct_posts']):.1f}%), and "
        f"dinosaurs or museums appear in {int(dinosaurs['count']):,} ({float(dinosaurs['pct_posts']):.1f}%). These are "
        f"topic-presence counts. They do not establish positive opinion, severity, or impact. They support a directional "
        f"content test using Chinese-language discovery cards.",
    )
    add_inline_section(
        doc,
        "Chinese-language Google validation",
        f"The primary model includes {cn_google_manifest['metrics']['n_reviews']:,} Chinese-language Google reviews "
        f"from {cn_google_manifest['metrics']['unique_pois']:,} POIs, with mean rating "
        f"{cn_google_manifest['metrics']['mean_review_rating']:.2f} stars. Automated SnowNLP sentiment was rejected as "
        f"a headline result: it labeled {validation['snownlp_negative_n']:,} reviews negative, but "
        f"{validation['snownlp_negative_rated_4_or_5_n']:,} "
        f"({pct(validation['snownlp_negative_rated_4_or_5_share'], 0)}) had 4-5 stars and only "
        f"{validation['snownlp_negative_rated_2_or_less_n']:,} had 2 stars or fewer. Chinese Google results therefore "
        f"use reviewer-provided star ratings.",
    )
    add_inline_section(
        doc,
        "POI-level opportunities",
        f"The exploratory index identifies {poi_metrics['n_fix_it']:,} fix-it candidates "
        f"({poi_metrics['n_fix_it_fukui']:,} in Fukui), {poi_metrics['n_promote_it']:,} promote-it candidates "
        f"({poi_metrics['n_promote_it_fukui']:,} in Fukui), and {poi_metrics['n_crowding_hotspots']:,} crowding hotspots. "
        f"No site met the strict promote-it rule. Fukui examples are {promo1['poi_name']} (n={int(promo1['n_reviews'])}, "
        f"{pct(promo1['positive_share'])} 4-5-star share, 95% Wilson CI {pct(promo1['positive_share_ci_low'])} to "
        f"{pct(promo1['positive_share_ci_high'])}) and {promo2['poi_name']} (n={int(promo2['n_reviews'])}, "
        f"{pct(promo2['positive_share'])}, 95% CI {pct(promo2['positive_share_ci_low'])} to "
        f"{pct(promo2['positive_share_ci_high'])}). Review volume reflects the collection window and cap. It is not "
        f"a visitor count.",
    )
    add_inline_section(
        doc,
        "Limits",
        "Reviews are observational and self-selected. Associations may reflect salience, writing style, POI mix, or "
        "unmeasured factors. POI clustering is not modeled in the Firth estimates, so uncertainty may be understated. "
        "Rare tags produce wide intervals. Chinese-language Google tags and Xiaohongshu topic codes record topic "
        "presence and may not identify complaint polarity. Evidence from Google reviews and Xiaohongshu remains separate.",
    )
    add_inline_section(
        doc,
        "Next step",
        f"Pre-register a visitor-session A/B test of the {priority1['solution_label_en']}. Randomize exposure, log "
        f"interactions, and measure clicks, saves, and itinerary additions. Specify the sample-size plan, primary "
        f"outcome, analysis model, confidence interval, and stopping rule before data collection.",
    )
    add_inline_section(
        doc,
        "Reproducibility",
        "Aggregate outputs preserve the command, input path, SHA256 hash, generation time, denominator, and caveats. "
        "Primary analysis uses reviewed codebooks and star ratings. Model sentiment remains a secondary check.",
        size=10.8,
    )
    add_source_line(
        doc,
        [
            PRIORITY_CSV,
            PRIORITY_MANIFEST,
            XHS_TOPICS_CSV,
            XHS_TOPICS_MANIFEST,
            CHINESE_GOOGLE_MANIFEST,
            POI_CSV,
            POI_MANIFEST,
        ],
    )

    doc.save(OUT_HANDOUT)


def add_question(doc: Document, question: str, answer: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(1)
    q = p.add_run("Q: " + question)
    q.bold = True
    q.font.color.rgb = RGBColor(0x16, 0x34, 0x5E)
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.35)
    p.paragraph_format.space_after = Pt(4)
    p.add_run("A: " + answer)


def build_qa() -> None:
    def deck_notes(path: Path, expected: int) -> list[str]:
        prs = Presentation(path)
        ns = [s.notes_slide.notes_text_frame.text.strip() for s in prs.slides]
        if len(ns) != expected:
            raise ValueError(f"{path.name}: expected {expected} notes, found {len(ns)}")
        return ns

    live = deck_notes(PPTX, 13)
    intro = deck_notes(INTRO_PPTX, 2)
    methods = deck_notes(METHODS_PPTX, 3)
    results = deck_notes(RESULTS_PPTX, 5)
    # Assembled presented order (13 slides). The simplified decks supply most
    # slides; the title, the two-priorities slide, and the close stay live.
    notes = [
        live[0],     # 1  Title and scope
        intro[0],    # 2  Research question
        intro[1],    # 3  Our approach: small nudges
        methods[0],  # 4  The data
        methods[1],  # 5  Turning review text into signals
        methods[2],  # 6  How we tested it (statistical model + rules, merged)
        results[0],  # 7  Which problems predict low ratings
        live[7],     # 8  Two priority information nudges
        results[1],  # 9  A separate idea from Chinese posts
        results[2],  # 10 Where to act: fix-it and promote-it sites
        results[3],  # 11 Rank common nudges by impact, then ease
        results[4],  # 12 What this can and cannot claim
        live[12],    # 13 Next step: test the first nudge
    ]

    titles = [
        "Title and scope",
        "Research question",
        "Our approach: small nudges",
        "The data",
        "Turning review text into signals",
        "How we tested it",
        "Which problems predict low ratings",
        "Two priority information nudges",
        "A separate idea from Chinese posts",
        "Where to act: fix-it and promote-it sites",
        "Rank common nudges by impact, then ease",
        "What this can and cannot claim",
        "Next step: test the first nudge",
    ]
    sources: list[list[Path]] = [
        [POI_CSV, POI_MANIFEST],                # 1  Title
        [ASPECT_MANIFEST, POI_MANIFEST],        # 2  Research question
        [ASPECT_MANIFEST],                      # 3  Our approach
        [ASPECT_MANIFEST, POI_MANIFEST],        # 4  The data
        [ASPECT_CSV, ASPECT_MANIFEST],          # 5  Turning review text into signals
        [ASPECT_CSV, ASPECT_MANIFEST],          # 6  How we tested it (model + rules)
        [ASPECT_CSV, ASPECT_MANIFEST],          # 7  Which problems predict low ratings
        [ASPECT_CSV, ASPECT_MANIFEST],          # 8  Two priority information nudges
        [XHS_TOPICS_CSV, XHS_TOPICS_MANIFEST],  # 9  A separate idea from Chinese posts
        [POI_CSV, POI_MANIFEST],                # 10 Where to act: fix-it and promote-it
        [PRIORITY_CSV, PRIORITY_MANIFEST],      # 11 Rank common nudges
        [ASPECT_MANIFEST, POI_MANIFEST],        # 12 What this can and cannot claim
        [PRIORITY_CSV, PRIORITY_MANIFEST],      # 13 Next step
    ]
    qa: list[list[tuple[str, str]]] = [
        [
            (
                "Why use the word Hokuriku when the project is Fukui-first?",
                f"The Google corpus covers {poi_metrics['n_pois_by_prefecture']['Fukui']} Fukui, "
                f"{poi_metrics['n_pois_by_prefecture']['Ishikawa']} Ishikawa, and "
                f"{poi_metrics['n_pois_by_prefecture']['Toyama']} Toyama sites. Pooling improves statistical support. "
                "Recommendations remain Fukui-first, and regional results provide context.",
            ),
            (
                "What is a nudge in this project?",
                "A small change to information or the choice environment that preserves options. Examples include "
                "showing hours, realistic duration, route order, or alternative sites before a visitor decides.",
            ),
        ],
        [
            (
                "Why use online reviews?",
                "Reviews contain naturally occurring descriptions of pain points and attractions. They support "
                "hypothesis generation. Self-selection prevents population-level or causal interpretation.",
            ),
            (
                "What does fix-it versus promote-it mean?",
                "Fix-it identifies sites with elevated pain-point evidence. Promote-it identifies lower-volume sites "
                "with high positive star-rating shares. Both are exploratory follow-up categories.",
            ),
        ],
        [
            (
                "Why information nudges instead of bigger fixes?",
                "Information nudges are cheap to build and quick to test, and they target problems a visitor can act "
                "on before the trip. Operator fixes such as price or cleanliness require the site operator, so they "
                "are flagged separately rather than nudged.",
            ),
            (
                "What are the two levers on this slide?",
                "Information provision gives clearer pre-visit information. Demand redistribution guides visitors "
                "toward quieter, under-visited but well-rated places. Both are low-cost prompts that preserve "
                "visitor choice.",
            ),
        ],
        [
            (
                f"Why are there {total_n:,} reviews but only {model_n:,} model rows?",
                f"The primary model includes Japanese ({languages['japanese']:,}), English "
                f"({languages['english']:,}), and Chinese-language ({languages['chinese']:,}) reviews. The remaining "
                f"{metrics['dropped_unsupported_language_rows']:,} other-language or undetected rows are excluded.",
            ),
            (
                "Do language groups represent nationality?",
                "No. They represent the language detected in the review text. Nationality is neither observed nor inferred.",
            ),
            (
                "What are the two Chinese-language sources?",
                f"The rating model uses {languages['chinese']:,} Chinese-language Google reviews. A separate set of "
                f"{int(scenic['denominator_posts']):,} Xiaohongshu posts has no star ratings and is used only for "
                "directional topic evidence.",
            ),
        ],
        [
            (
                "How were review aspects identified?",
                "Human-reviewed multilingual keyword codebooks tag explicit aspect terms. Matched-term evidence and "
                "review decisions are retained for audit. Tags show that a topic was mentioned.",
            ),
            (
                f"Why define low rating as {low_rating_cutoff} stars or fewer?",
                f"The analysis manifest defines low rating as {aspect_manifest['filters']['low_rating_definition']}. "
                f"This produces {low_n:,} low-rated rows among {model_n:,} modeled reviews.",
            ),
            (
                "Why use stars instead of model sentiment?",
                "Google stars provide the same 1-5 outcome across languages. VADER, oseti, and SnowNLP do not share a "
                "measurement scale and remain secondary within-language checks.",
            ),
        ],
        [
            (
                "Why Firth logistic regression?",
                "Rare tags can create small-sample bias or separation in ordinary logistic regression. Firth bias "
                "reduction yields finite, more stable estimates under sparsity.",
            ),
            (
                "What variables does the model adjust for?",
                "Review text length, review language, and prefecture. Adjustment reduces confounding from longer texts "
                "matching more keywords and from regional or language-group composition.",
            ),
            (
                "What does an odds ratio mean?",
                "It compares odds after adjustment. OR 4.77 means 4.77 times the odds, not 4.77 times the probability. "
                "Association does not establish causal impact.",
            ),
            (
                "Why use Benjamini-Hochberg FDR?",
                f"The primary analysis fits {metrics['A_primary_models_fit']} aspect models. Multiple testing increases "
                "false-positive risk. BH-FDR controls the expected false-discovery proportion among rejected tests.",
            ),
            (
                "What is required for a ranked opportunity?",
                f"A harmful direction, BH-FDR significance, at least {min_mentions} pooled mentions, and a configured "
                "action type. The action mapping separates information nudges from operator fixes.",
            ),
            (
                "How robust are the results to ordinary logistic regression?",
                f"For {sanity_aspects} selected aspects, the maximum absolute difference in log odds ratio was "
                f"{max(v['log_abs_diff'] for v in metrics['plain_vs_firth_sanity'].values()):.3f}. This is a limited "
                "sanity check, not a full robustness analysis.",
            ),
        ],
        [
            (
                "Why trust the itinerary result with only 29 mentions?",
                f"The adjusted OR is {itinerary['odds_ratio']:.2f}, 95% CI {itinerary['or_ci_low']:.2f} to "
                f"{itinerary['or_ci_high']:.2f}, and BH-FDR p={fmt_p(itinerary['p_value_bh_fdr'])}. Firth addresses "
                "sparsity, while the wide interval shows remaining uncertainty.",
            ),
            (
                "Is wayfinding statistically significant?",
                f"Its BH-FDR p-value is {fmt_p(wayfinding['p_value_bh_fdr'])}, but the Wald 95% CI is "
                f"{wayfinding['or_ci_low']:.2f} to {wayfinding['or_ci_high']:.2f}, which includes 1. Different inferential "
                "summaries sit near the threshold. The honest label is preliminary or borderline.",
            ),
            (
                "Why are explicit pain-point counts low?",
                "The codebooks identify explicit problem-language matches rather than every discussion of a topic. "
                "Complaint behavior also varies by language, platform, and writing style.",
            ),
        ],
        [
            (
                "Why prioritize opening hours and itinerary information?",
                f"Opening hours has {int(opening['n_positive'])} mentions, OR {opening['odds_ratio']:.2f}, and 95% CI "
                f"{opening['or_ci_low']:.2f} to {opening['or_ci_high']:.2f}. Itinerary/time has "
                f"{int(itinerary['n_positive'])} mentions, OR {itinerary['odds_ratio']:.2f}, and 95% CI "
                f"{itinerary['or_ci_low']:.2f} to {itinerary['or_ci_high']:.2f}. Both are actionable before a trip.",
            ),
            (
                "Why exclude price even though its odds ratio is larger?",
                f"Price has OR {price['odds_ratio']:.2f}, but changing price is an operator decision. The study focuses "
                "on low-cost information nudges, so price is flagged for operational follow-up.",
            ),
        ],
        [
            (
                "Why keep Xiaohongshu separate?",
                "The posts have no star-rating outcome and come from one pre-trip social platform. They cannot enter the "
                "Google rating model. Their role is directional content discovery.",
            ),
            (
                "What do the Chinese topic counts mean?",
                f"Scenic nature appears in {int(scenic['count'])} of {int(scenic['denominator_posts'])} posts "
                f"({float(scenic['pct_posts']):.1f}%). Dinosaurs or museums appear in {int(dinosaurs['count'])} "
                f"({float(dinosaurs['pct_posts']):.1f}%). A tag means topic presence, not sentiment or importance.",
            ),
            (
                "Why not report SnowNLP sentiment?",
                f"Validation against Chinese-language Google stars failed. SnowNLP labeled "
                f"{validation['snownlp_negative_n']} reviews negative; {validation['snownlp_negative_rated_4_or_5_n']} "
                f"({pct(validation['snownlp_negative_rated_4_or_5_share'], 0)}) had 4-5 stars, while only "
                f"{validation['snownlp_negative_rated_2_or_less_n']} had 2 stars or fewer. Stars are used instead.",
            ),
        ],
        [
            (
                "Are promote-it sites proven tourism opportunities?",
                f"No. The index identifies {poi_metrics['n_promote_it']} exploratory candidates and "
                f"{poi_metrics['n_promote_it_strict']} strict candidates. Point-estimate satisfaction and collection "
                "volume generate follow-up candidates.",
            ),
            (
                "Does low review volume mean low visitor volume?",
                "No. Review collection was capped and reflects a collection window and POI mix. Volume categories are "
                "relative to this corpus.",
            ),
            (
                "How uncertain are the Fukui examples?",
                f"{promo1['poi_name']} has n={int(promo1['n_reviews'])}, {pct(promo1['positive_share'])} positive, and "
                f"95% Wilson CI {pct(promo1['positive_share_ci_low'])} to {pct(promo1['positive_share_ci_high'])}. "
                f"{promo2['poi_name']} has n={int(promo2['n_reviews'])}, {pct(promo2['positive_share'])}, and CI "
                f"{pct(promo2['positive_share_ci_low'])} to {pct(promo2['positive_share_ci_high'])}.",
            ),
        ],
        [
            (
                "Why is the visit-readiness card ranked first?",
                f"It combines the opening-hours and itinerary rating associations with an Easy implementation tier in "
                "the pre-specified mapping. It also has reviewed support across all three language/source groups.",
            ),
            (
                "Does support across all three groups mean equal evidence?",
                "No. English and Japanese Google evidence forms the confirmatory backbone. Chinese-language XHS topic "
                "evidence is directional. Source denominators and instruments remain separate.",
            ),
            (
                "Why call crowding preliminary if BH-FDR clears the threshold?",
                f"Waiting/crowding has OR {waiting['odds_ratio']:.2f}, BH-FDR p={fmt_p(waiting['p_value_bh_fdr'])}, and "
                f"95% CI {waiting['or_ci_low']:.2f} to {waiting['or_ci_high']:.2f}. The interval includes 1, so the "
                "association is treated cautiously.",
            ),
        ],
        [
            (
                "What is the largest threat to validity?",
                "Observational self-selection and unmodeled POI nesting. Reviews from the same site are not independent, "
                "and the row-level Firth model may understate uncertainty.",
            ),
            (
                "Can the findings generalize to all tourists?",
                "No population-representative claim is made. Findings concern review writers and Xiaohongshu users in "
                "the observed collection window and POI mix.",
            ),
        ],
        [
            (
                "What is the first experiment?",
                f"A visitor-session A/B test of the {priority1['solution_label_en']}. The treatment displays hours, "
                "duration, transfers, booking needs, and backups before itinerary commitment.",
            ),
            (
                "What outcomes should be measured?",
                "Primary behavioral outcomes are card clicks, saves, and itinerary additions. Low-rating share is a "
                "longer-term secondary outcome requiring suitable follow-up or linkage.",
            ),
            (
                "What is needed before launch?",
                "A pre-registered hypothesis, sample-size plan, randomization procedure, primary outcome, analysis model, "
                "confidence interval, stopping rule, and data-quality checks.",
            ),
        ],
    ]

    doc = Document()
    configure_a4(doc, 1.7)
    set_default_styles(doc, 10.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Seminar Speaker Script and Statistical Q&A")
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0x16, 0x34, 0x5E)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(8)
    r = p.add_run("Green ANDREW and Xu ZILIN  |  June 30, 2026  |  10-minute presentation")
    r.bold = True

    add_inline_section(
        doc,
        "Delivery plan",
        "Green ANDREW leads slides 1-3, 5-8, and 11-13. Xu ZILIN presents slides 4, 9, and 10. Speak slowly, pause "
        "after key numbers, and read only the headline statistics. Detailed methods and caveats below are for professor Q&A.",
    )

    for index, (title_text, note, questions, slide_sources) in enumerate(
        zip(titles, notes, qa, sources), start=1
    ):
        doc.add_heading(f"Slide {index}: {title_text}", level=1)
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        set_cell = p.add_run("Spoken script\n")
        set_cell.bold = True
        set_cell.font.color.rgb = RGBColor(0x16, 0x34, 0x5E)
        p.add_run(note)
        doc.add_heading("Likely questions", level=2)
        for question, answer in questions:
            add_question(doc, question, answer)
        add_source_line(doc, slide_sources)

    doc.add_heading("Core defense rules", level=1)
    for text in (
        "Say association, not impact or cause.",
        "An odds ratio multiplies odds. It is not a probability ratio.",
        "Wayfinding and crowding are preliminary because their 95% intervals include 1.",
        "Chinese friction and topic tags record topic presence. They do not establish polarity or severity.",
        "Xiaohongshu is a separate directional source without star ratings.",
        "Language groups describe text language, not nationality.",
        "POI classifications and solution ranks identify follow-up experiments. They do not estimate intervention effectiveness.",
    ):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(text)

    doc.add_heading("Primary reproducibility sources", level=1)
    for path in (
        ASPECT_CSV,
        ASPECT_MANIFEST,
        POI_CSV,
        POI_MANIFEST,
        PRIORITY_CSV,
        PRIORITY_MANIFEST,
        XHS_TOPICS_CSV,
        XHS_TOPICS_MANIFEST,
        CHINESE_GOOGLE_MANIFEST,
    ):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(path.relative_to(ROOT).as_posix())

    doc.save(OUT_QA)


def main() -> int:
    build_qa()
    build_handout()
    print(f"wrote {OUT_QA}")
    print(f"wrote {OUT_HANDOUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
