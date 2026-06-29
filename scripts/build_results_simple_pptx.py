#!/usr/bin/env python3
"""Build NUDGE-Results-Simple.pptx: de-jargoned drop-in replacements for three
RESULTS/DISCUSSION slides, to be hand-inserted into the live Google Drive deck.

Replaces live slides 7 (results table), 9 (Xiaohongshu), and 12 (can/cannot
claim). Principle for an engineering-lab audience: keep every number where we
SHOW evidence (odds ratios, 95% CI, BH-FDR p stay on the slide), but pair each
number with a plain-language reading, and de-jargon the prose that frames it.
Nothing here removes a statistic; it only adds an intuitive translation.

Provenance: every number is pulled live through the SAME deck getters as the
main deck, reused via build_nudge_pptx. The existing generation is untouched.

Slide numbering reflects the post-merge deck: after the 3 simplified methods
slides replace the old 4, the deck is 12 slides, and these three land at
positions 6, 8, and 11.

Run:  .venv/bin/python3 scripts/build_results_simple_pptx.py
"""
from __future__ import annotations

from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import build_nudge_pptx as bp

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "NUDGE-Results-Simple.pptx"
LADDER_SVG = ROOT / "docs/statistical_test_figures/figure_evidence_ladder.svg"
DECK_TOTAL = 12


def rasterize(svg_path: Path) -> Path:
    if not svg_path.exists():
        raise SystemExit(f"missing figure: {svg_path}")
    png = bp.SCRATCH / f"results_simple_{svg_path.stem}.png"
    cairosvg.svg2png(url=str(svg_path), write_to=str(png), output_width=1600)
    return png


LADDER_PNG = rasterize(LADDER_SVG)


def plain_x(or_str: str) -> str:
    """Odds ratio -> rounded 'about N x' plain reading."""
    return f"{float(or_str):.1f}"


def dejargon_en(t: str) -> str:
    """De-jargon a LIVE evidence string by substitution (no number is retyped)."""
    return (str(t)
            .replace("through FDR", "through the strict test")
            .replace("clears FDR", "passes the strict test")
            .replace("clear FDR", "pass the strict test")
            .replace("XHS", "the Chinese posts")
            .replace("Google models", "the rating model"))


def dejargon_ja(t: str) -> str:
    return (str(t)
            .replace("FDRを通過しない", "厳しい基準を通過しない")
            .replace("FDRを通過", "厳しい基準を通過")
            .replace("XHS", "中国語投稿")
            .replace("Googleモデル", "評価モデル"))


def picture_fit(slide, png: Path, x, y, max_w, max_h):
    from PIL import Image

    iw, ih = Image.open(png).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w, h = max_w, int(max_w / ar)
    else:
        h, w = max_h, int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    slide.shapes.add_picture(str(png), Emu(int(px)), Emu(int(py)), Emu(int(w)), Emu(int(h)))


def foot(slide, num: int):
    tb, tf = bp.textbox(slide, Inches(11.0), Inches(7.04), Inches(1.83), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{num} / {DECK_TOTAL}"
    bp._style_run(r, 10, bp.GREY)


def build():
    prs = Presentation()
    prs.slide_width = bp.EMU_W
    prs.slide_height = bp.EMU_H
    blank = prs.slide_layouts[6]

    def new():
        s = prs.slides.add_slide(blank)
        bp._set_bg(s, bp.WHITE)
        return s

    # ===== SLIDE 1 (replaces live slide 7): results table ====================
    s = new()
    bp.label(s, "III.  RESULTS")
    bp.title(s, "Which problems predict low ratings", "どの不満が低評価を予測するか")

    rows = [
        ("Opening hours / availability", "開館時間・営業状況", bp.OPEN_PREV,
         f"{bp.OPEN_OR} ({bp.OPEN_CIL}–{bp.OPEN_CIH})", bp.OPEN_PV, "Yes", "はい"),
        ("Itinerary fit / time-cost", "行程適合・所要時間", bp.TIME_PREV,
         f"{bp.TIME_OR} ({bp.TIME_CIL}–{bp.TIME_CIH})", bp.TIME_PV, "Yes", "はい"),
        ("Wayfinding / signage", "道案内・表示", bp.SIGN_PREV,
         f"{bp.SIGN_OR} ({bp.SIGN_CIL}–{bp.SIGN_CIH})", bp.SIGN_PV, "Yes", "はい"),
    ]
    headers = [
        ("Pain point", "不満点"), ("Seen in", "出現率"),
        ("Adjusted OR (95% CI)", "調整OR(95%CI)"), ("BH-FDR p", "BH-FDR p"),
        ("Nudge-able", "ナッジ可能"),
    ]
    col_w = [Inches(2.5), Inches(0.9), Inches(2.15), Inches(1.1), Inches(1.0)]
    tbl_w = sum(col_w, Inches(0))
    gframe = s.shapes.add_table(len(rows) + 1, len(headers), bp.MX, Inches(2.0),
                                tbl_w, Inches(2.2))
    table = gframe.table
    table.first_row = False
    table.horz_banding = False
    for j, w in enumerate(col_w):
        table.columns[j].width = w
    for j, (en, jp) in enumerate(headers):
        bp._cell(table.cell(0, j), en, 12, bp.WHITE, bold=True,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=bp.NAVY, jp=jp)
    for i, (en, jp, prev, orci, p, nud_en, nud_jp) in enumerate(rows, start=1):
        rowfill = bp.WHITE if i % 2 else bp.CARD_BG
        bp._cell(table.cell(i, 0), en, 12, bp.INK, bold=True, fill=rowfill, jp=jp)
        bp._cell(table.cell(i, 1), prev, 12, bp.INK, align=PP_ALIGN.CENTER, fill=rowfill)
        bp._cell(table.cell(i, 2), orci, 12, bp.NAVY, bold=True, align=PP_ALIGN.CENTER, fill=rowfill)
        bp._cell(table.cell(i, 3), p, 12, bp.INK, align=PP_ALIGN.CENTER, fill=rowfill)
        bp._cell(table.cell(i, 4), nud_en, 12, bp.BLUE, bold=True, align=PP_ALIGN.CENTER,
                 fill=rowfill, jp=nud_jp)

    picture_fit(s, bp.PNG["nudge_aspect_fig"], Inches(8.5), Inches(2.0), Inches(4.3), Inches(2.7))

    bp.infobox(s, bp.MX, Inches(4.55), Inches(7.85), Inches(1.35),
               f"In plain terms: a review that mentions opening hours is about {plain_x(bp.OPEN_OR)}x more likely "
               f"to be low-rated. Itinerary problems: about {plain_x(bp.TIME_OR)}x. Wayfinding: about "
               f"{plain_x(bp.SIGN_OR)}x (borderline — its range just touches 1).",
               f"平たく言えば：開館時間に触れた口コミは低評価が約{plain_x(bp.OPEN_OR)}倍になりやすい。"
               f"行程の問題は約{plain_x(bp.TIME_OR)}倍、道案内は約{plain_x(bp.SIGN_OR)}倍"
               f"(境界的—範囲がちょうど1に接する)。")
    bp.caption(s, bp.MX, Inches(6.05), Inches(12.0),
               "Adjusted OR = odds ratio after accounting for review length, language, and prefecture. "
               f"Higher means more linked to low ratings. Price (OR {bp.PRICE_OR}) and cleanliness (OR {bp.CLEAN_OR}) "
               "matter too, but need an operator fix.",
               "調整OR=口コミの長さ・言語・県をそろえた後のオッズ比。大きいほど低評価と強く関連。"
               f"価格(OR {bp.PRICE_OR})や清潔さ(OR {bp.CLEAN_OR})も関連するが、事業者の対応が必要。")
    foot(s, 6)
    bp.notes(s,
             f"Andrew (55 sec): Three nudge-able pain points are significant after FDR correction. The odds ratios "
             f"are {bp.OPEN_OR} for opening hours, {bp.TIME_OR} for itinerary fit, and {bp.SIGN_OR} for wayfinding. "
             f"In plain terms that is about {plain_x(bp.OPEN_OR)}, {plain_x(bp.TIME_OR)}, and {plain_x(bp.SIGN_OR)} "
             "times the odds of a low rating when the problem is mentioned. Wayfinding is borderline: its 95% "
             f"interval ({bp.SIGN_CIL} to {bp.SIGN_CIH}) just includes 1. Price (OR {bp.PRICE_OR}) and cleanliness "
             f"(OR {bp.CLEAN_OR}) are also linked but require operator action, not information.")

    # ===== SLIDE 2 (replaces live slide 9): Xiaohongshu hypothesis ===========
    s = new()
    bp.label(s, "III.  RESULTS")
    bp.title(s, "A separate idea from Chinese posts", "中国語の投稿から得た、別の手がかり")

    _, c1 = bp.card(s, bp.MX, Inches(2.05), Inches(6.0), Inches(3.7))
    bp.en_jp(c1, "What the Chinese posts hint at", "中国語の投稿が示唆すること",
             en_size=17, jp_size=12.5, en_color=bp.NAVY, bold=True, first=True, space_after=4)
    bp.en_jp(c1, "(weaker evidence)", "(弱い根拠)", en_size=13, jp_size=11, en_color=bp.GREY, space_after=12)
    bp.en_jp(c1, f"{bp.CN_ROWS} Chinese-language Fukui posts from Xiaohongshu. There is no star rating here, so this is a hint, not a test.",
             f"福井に関する中国語の投稿{bp.CN_ROWS}件(小紅書)。星評価がないため、これは「検証」ではなく「手がかり」。",
             en_size=14, jp_size=11, en_color=bp.INK, space_after=12)
    bp.en_jp(c1, "Sentiment was read with SnowNLP (a Chinese-text sentiment tool), interpreted only within this source.",
             "感情はSnowNLP(中国語テキストの感情分析ツール)で判定し、このソース内だけで解釈する。",
             en_size=14, jp_size=11, en_color=bp.INK, space_after=0)

    _, c2 = bp.card(s, Inches(6.85), Inches(2.05), Inches(5.75), Inches(3.7))
    bp.en_jp(c2, "Signals worth testing", "検証する価値のあるシグナル",
             en_size=17, jp_size=12.5, en_color=bp.NAVY, bold=True, first=True, space_after=12)
    bp.en_jp(c2, f"Dinosaur / museum: {bp.DINO_POS} of {bp.DINO_N} posts positive ({bp.DINO_PCT}), versus {bp.DINO_OTHER_PCT} without the tag ({bp.DINO_FDR}).",
             f"恐竜・博物館：{bp.DINO_N}件中{bp.DINO_POS}件がポジティブ({bp.DINO_PCT})。タグなしは{bp.DINO_OTHER_PCT}({bp.DINO_FDR})。",
             en_size=14, jp_size=11, en_color=bp.INK, space_after=12)
    bp.en_jp(c2, f"Scenic nature: {bp.SCENIC_POS} of {bp.SCENIC_N} positive ({bp.SCENIC_PCT}), versus {bp.SCENIC_OTHER_PCT} ({bp.SCENIC_FDR}).",
             f"自然景観：{bp.SCENIC_N}件中{bp.SCENIC_POS}件がポジティブ({bp.SCENIC_PCT})。タグなしは{bp.SCENIC_OTHER_PCT}({bp.SCENIC_FDR})。",
             en_size=14, jp_size=11, en_color=bp.INK, space_after=0)

    bp.infobox(s, bp.MX, Inches(6.0), Inches(12.1), Inches(1.0),
               "Candidate nudge: test a Chinese-language discovery card featuring dinosaurs / museums and scenery. "
               "Hint only — one platform, a secondary sentiment tool, no rating model, no causal claim.",
               "候補ナッジ：恐竜・博物館と自然景観を前面に出した中国語の発見カードを検証する。"
               "手がかりに限定—単一プラットフォーム、副次的な感情ツール、星評価モデルなし、因果主張なし。")
    foot(s, 8)
    bp.notes(s,
             f"Lynn (55 sec): Separately, we looked at {bp.CN_ROWS} Chinese-language posts about Fukui from "
             "Xiaohongshu. These have no star rating, so we treat them as a hint, not a test. Posts tagged with "
             f"dinosaurs or museums were positive {bp.DINO_PCT} of the time, versus {bp.DINO_OTHER_PCT} without the "
             f"tag ({bp.DINO_FDR} after BH-FDR). Scenic-nature posts were {bp.SCENIC_PCT} positive versus "
             f"{bp.SCENIC_OTHER_PCT} ({bp.SCENIC_FDR}). This suggests a Chinese-language discovery card to A/B test "
             "next. It is hypothesis-generating only: one platform, a secondary sentiment tool, no causal claim.")

    # ===== SLIDE 3 (replaces live slide 10): where to act ====================
    s = new()
    bp.label(s, "III.  RESULTS")
    bp.title(s, "Where to act: fix-it and promote-it sites", "どこで動くか：改善型と推奨型のスポット")
    bp.stat_callout(s, bp.MX, Inches(2.0), Inches(2.3), bp.FIX_COUNT, "Fix-it sites", "改善型")
    bp.stat_callout(s, Inches(3.05), Inches(2.0), Inches(2.3), bp.PROMOTE_COUNT, "Promote-it sites", "推奨型")
    bp.stat_callout(s, Inches(5.5), Inches(2.0), Inches(2.3), bp.CROWD_COUNT, "Crowding hotspots", "混雑ホット")
    _, tf = bp.textbox(s, bp.MX, Inches(3.95), Inches(7.7), Inches(3.0))
    bp.en_jp(tf,
             f"Fix-it sites show strong pain points; promote-it sites have a high share of 4-5 star reviews "
             f"({bp.FIX_COUNT} fix-it, {bp.FIX_FUKUI} in Fukui; {bp.PROMOTE_COUNT} promote-it, {bp.PROMOTE_FUKUI} in Fukui; {bp.CROWD_COUNT} crowding hotspots).",
             f"改善型は不満点が強いスポット、推奨型は4〜5つ星の割合が高いスポット"
             f"(改善型{bp.FIX_COUNT}・うち福井{bp.FIX_FUKUI}、推奨型{bp.PROMOTE_COUNT}・うち福井{bp.PROMOTE_FUKUI}、混雑ホット{bp.CROWD_COUNT})。",
             first=True, bullet=True, en_size=14, jp_size=11, space_after=14)
    bp.en_jp(tf,
             f"More reviews mean more trust: under {bp.LOW_VOL} reviews is thin, over {bp.HIGH_VOL} is solid.",
             f"口コミが多いほど信頼できる：{bp.LOW_VOL}件未満は少数、{bp.HIGH_VOL}件超は多数。",
             bullet=True, en_size=14, jp_size=11, space_after=14)
    bp.en_jp(tf,
             f"Top Fukui promote-it: {bp.PROMO1} {bp.PROMO1_SHARE} positive (95% range {bp.PROMO1_LOW}-{bp.PROMO1_HIGH}); "
             f"then {bp.PROMO2} {bp.PROMO2_SHARE} ({bp.PROMO2_LOW}-{bp.PROMO2_HIGH}).",
             f"福井の推奨型トップ：{bp.PROMO1} {bp.PROMO1_SHARE}(95%範囲 {bp.PROMO1_LOW}〜{bp.PROMO1_HIGH})、"
             f"次に {bp.PROMO2} {bp.PROMO2_SHARE}({bp.PROMO2_LOW}〜{bp.PROMO2_HIGH})。",
             bullet=True, en_size=14, jp_size=11, space_after=0)
    picture_fit(s, bp.PNG["nudge_poi_fig"], Inches(8.5), Inches(2.0), Inches(4.3), Inches(4.1))
    foot(s, 9)
    bp.notes(s,
             f"Lynn (55 sec): We sort sites into three groups. Fix-it sites ({bp.FIX_COUNT}, {bp.FIX_FUKUI} in Fukui) "
             f"have strong pain points. Promote-it sites ({bp.PROMOTE_COUNT}, {bp.PROMOTE_FUKUI} in Fukui) have a high "
             f"share of 4-5 star reviews. We also flag {bp.CROWD_COUNT} crowding hotspots. We trust sites with more "
             f"reviews: under {bp.LOW_VOL} is thin, over {bp.HIGH_VOL} is solid. Shares use Wilson 95% confidence "
             f"intervals. The top Fukui promote-it site is {bp.PROMO1} at {bp.PROMO1_SHARE} positive (Wilson 95% CI "
             f"{bp.PROMO1_LOW} to {bp.PROMO1_HIGH}), then {bp.PROMO2} at {bp.PROMO2_SHARE}.")

    # ===== SLIDE 4 (replaces live slide 11): rank common nudges ==============
    s = new()
    bp.label(s, "III.  RESULTS")
    bp.title(s, "Rank common nudges by impact, then ease", "共通ナッジをインパクト、次に実装容易性で順位づける")
    _, itf = bp.textbox(s, bp.MX, Inches(1.92), Inches(12.1), Inches(0.6))
    bp.en_jp(itf,
             "Each idea has reviewed support from English, Japanese, and Chinese-language reviews. Evidence types stay separate.",
             "各施策には英語・日本語・中国語のレビュー済み根拠がある。エビデンス種別は統合しない。",
             first=True, en_size=13, jp_size=10.5, en_color=bp.GREY, space_after=0)
    headers = [("Rank", "順位"), ("Common solution", "共通施策"),
               ("Why it ranks here", "順位の理由"), ("Next-semester test", "来学期の実験")]
    col_w = [Inches(0.7), Inches(3.0), Inches(4.7), Inches(3.7)]
    gframe = s.shapes.add_table(len(bp.PRIORITIES) + 1, len(headers), bp.MX, Inches(2.65),
                                sum(col_w, Inches(0)), Inches(3.7))
    table = gframe.table
    table.first_row = False
    table.horz_banding = False
    for j, w in enumerate(col_w):
        table.columns[j].width = w
    for j, (en, jp) in enumerate(headers):
        bp._cell(table.cell(0, j), en, 12, bp.WHITE, bold=True,
                 align=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT, fill=bp.NAVY, jp=jp)
    for i, p in enumerate(bp.PRIORITIES, start=1):
        rowfill = bp.WHITE if i % 2 else bp.CARD_BG
        bp._cell(table.cell(i, 0), p["rank"], 16, bp.NAVY, bold=True, align=PP_ALIGN.CENTER, fill=rowfill)
        bp._cell(table.cell(i, 1), p["name_en"], 12, bp.INK, bold=True, fill=rowfill, jp=p["name_ja"])
        bp._cell(table.cell(i, 2),
                 f"Impact: {p['impact']} · Build: {p['ease']}. {dejargon_en(p['summary_en'])}",
                 10.5, bp.INK, fill=rowfill,
                 jp=f"インパクト:{p['impact']} ・ 実装:{p['ease']}。{dejargon_ja(p['summary_ja'])}")
        bp._cell(table.cell(i, 3), dejargon_en(p["test_en"]), 10.5, bp.INK, fill=rowfill, jp=p["test_ja"])
    bp.infobox(s, bp.MX, Inches(6.5), Inches(12.1), Inches(0.66),
               "This is an opportunity ranking, not proof of effectiveness. Higher impact ranks first; easier to build breaks ties.",
               "これは機会の順位づけで、効果の証明ではない。インパクトが高い順、同点なら作りやすさで決める。")
    foot(s, 10)
    bp.notes(s,
             "Andrew (55 sec): We rank the cross-language solutions by impact first, then by how easy they are to "
             "build. Priority one is the multilingual visit-readiness card: all three languages flag planning "
             "problems, and the opening-hours and itinerary signals passed the strict test. Priority two is the "
             "localized discovery card, supported by the Chinese dinosaur and scenery signals. Priority three, "
             "off-peak prompts, shows crowding in all three languages; its link to low ratings passes the strict "
             "test but is borderline, and it ranks lower mainly because it is the hardest to build. This is an "
             "opportunity ranking, not proof of effectiveness.")

    # ===== SLIDE 5 (replaces live slide 12): what we can and cannot claim ====
    s = new()
    bp.label(s, "IV.  DISCUSSION")
    bp.title(s, "What this can and cannot claim", "主張できること・できないこと")

    _, tf = bp.textbox(s, bp.MX, Inches(2.05), Inches(7.0), Inches(4.6))
    bp.en_jp(tf, "This is exploratory: it points to experiments to run, not proof of cause and effect.",
             "これは探索的：実験の候補を示すもので、因果関係の証明ではない。",
             first=True, bullet=True, en_size=16, jp_size=12.5, space_after=18)
    bp.en_jp(tf, "The opportunity scores rank what to test next. They do not measure how well a fix works.",
             "機会スコアは「次に何を試すか」の順位づけ。効果の大きさを測るものではない。",
             bullet=True, en_size=16, jp_size=12.5, space_after=18)
    bp.en_jp(tf, "We treat each review on its own, so the model may be slightly overconfident — another reason we call this a ranking, not proof.",
             "各口コミを個別に扱うため、モデルはやや自信過剰かもしれない—だからこそ順位づけであり、証明ではない。",
             bullet=True, en_size=16, jp_size=12.5, space_after=18)
    bp.en_jp(tf, "Language groups describe the review's language, not the writer's nationality.",
             "言語グループは口コミの言語であり、執筆者の国籍ではない。",
             bullet=True, en_size=16, jp_size=12.5, space_after=0)

    picture_fit(s, LADDER_PNG, Inches(7.6), Inches(2.0), Inches(5.15), Inches(4.3))
    foot(s, 11)
    bp.notes(s,
             "Andrew (50 sec): We are clear about the limits. This study is exploratory and hypothesis-generating, "
             "not causal. The opportunity scores rank what to test next; they do not estimate how well an "
             "intervention works. We model each review on its own, so because reviews cluster within sites, the "
             "row-level estimates may understate uncertainty — which is exactly why we present a ranking rather "
             "than a proof. And language groups describe the review language, not the reviewer's nationality. "
             "The evidence ladder shows where we sit: above description, at association, below a causal claim.")

    prs.save(str(OUT_PPTX))
    return prs


def main() -> int:
    build()
    size = OUT_PPTX.stat().st_size
    print(f"wrote {OUT_PPTX} ({size:,} bytes); 5 slides (replaces live slides 7, 9, 10, 11, 12)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
