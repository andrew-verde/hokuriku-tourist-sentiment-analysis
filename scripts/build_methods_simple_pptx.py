#!/usr/bin/env python3
"""Build NUDGE-Methods-Simple.pptx: a 12th-grade-readable replacement for the
METHODS section only, to be hand-inserted into the live Google Drive deck.

Scope: replaces live slides 3-6 (Data, Text-into-signals, Statistical model,
Inference/decision rule) with THREE slides. Slides 5 and 6 are merged into one
plain "How we tested it" slide. The heavy statistics (Firth, odds ratios, 95%
CI, Benjamini-Hochberg FDR, the Firth-vs-logit sanity check) are stripped off
the slide FACE and kept in the speaker notes, where the spoken script lives; the
Seminar-Speaker-Script-and-QA.docx Q&A already carries the same detail.

Provenance: every displayed number is pulled live through the SAME deck getters
as the main deck. We import scripts/build_nudge_pptx.py (which resolves all
values via getters at import time) and reuse its styling helpers and constants.
Nothing in the existing generation is modified.

Calibration: detail level matches the reference student design-thinking deck
(stats are named, not explained); visual skin matches the live deck so the
slides drop in seamlessly. EN primary with a complete plain-language JP subtitle.

Run:  .venv/bin/python3 scripts/build_methods_simple_pptx.py
"""
from __future__ import annotations

from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Reuse the main PPTX module: importing it resolves every number through the
# deck getters (provenance) and gives us the shared style helpers. The heavy
# build() only runs under its own __main__ guard, so import is side-effect-safe.
import build_nudge_pptx as bp

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "NUDGE-Methods-Simple.pptx"
FUNNEL_SVG = ROOT / "docs/statistical_test_figures/figure_decision_rule_funnel_simple.svg"

# Merged deck end-state: replacing 4 method slides with 3 leaves a 12-slide deck.
DECK_TOTAL = 12


def rasterize(svg_path: Path) -> Path:
    if not svg_path.exists():
        raise SystemExit(f"missing figure: {svg_path}")
    png = bp.SCRATCH / f"methods_simple_{svg_path.stem}.png"
    cairosvg.svg2png(url=str(svg_path), write_to=str(png), output_width=1600)
    return png


FUNNEL_PNG = rasterize(FUNNEL_SVG)


def picture_fit(slide, png: Path, x, y, max_w, max_h):
    """Place a PNG centered inside a box, preserving aspect ratio."""
    from PIL import Image

    iw, ih = Image.open(png).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w, h = max_w, int(max_w / ar)
    else:
        h, w = max_h, int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    slide.shapes.add_picture(str(png), Emu(int(px)), Emu(int(py)), Emu(int(w)), Emu(int(h)))


def foot(slide, num: int):
    """Page number bottom-right, matching the live deck's 'N / total' style."""
    tb, tf = bp.textbox(slide, Inches(11.0), Inches(7.04), Inches(1.83), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{num} / {DECK_TOTAL}"
    bp._style_run(r, 10, bp.GREY)


def build():
    prs = Presentation()
    prs.slide_width = bp.EMU_W
    prs.slide_height = bp.EMU_H
    blank = prs.slide_layouts[6]

    def new():
        s = prs.slides.add_slide(blank)
        bp._set_bg(s, bp.WHITE)
        return s

    # ===== SLIDE A (replaces live slide 3): The data ==========================
    s = new()
    bp.label(s, "II.  METHODS")
    bp.title(s, "The data", "データ")

    _, c1 = bp.card(s, bp.MX, Inches(2.05), Inches(5.7), Inches(2.05))
    pe = c1.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Tagged reviews"
    bp._style_run(rr, 16, bp.NAVY, bold=True, font=bp.HEAD_FONT)
    pj = c1.add_paragraph()
    rj = pj.add_run()
    rj.text = "タグ付き口コミ"
    bp._style_run(rj, 11.5, bp.GREY, font=bp.JP_FONT)
    bp.en_jp(c1, f"{bp.TAGGED_ROWS} Google reviews", f"Google口コミ {bp.TAGGED_ROWS}件",
             en_size=17, jp_size=12, en_color=bp.NAVY, bold=True)
    bp.en_jp(c1, f"Japanese {bp.LANG_JP} · English {bp.LANG_EN} · Chinese {bp.LANG_CN}  (modeled)",
             f"日本語 {bp.LANG_JP} ・ 英語 {bp.LANG_EN} ・ 中国語 {bp.LANG_CN}(モデル対象)",
             en_size=14, jp_size=11, en_color=bp.INK, space_after=4)
    bp.en_jp(c1, f"Other-language {bp.LANG_OTHER} · too-short {bp.LANG_UNDETECTED}  (not modeled)",
             f"その他言語 {bp.LANG_OTHER} ・ 短すぎ {bp.LANG_UNDETECTED}(モデル対象外)",
             en_size=12.5, jp_size=10, en_color=bp.GREY, space_after=0)

    _, c2 = bp.card(s, bp.MX, Inches(4.3), Inches(5.7), Inches(1.95))
    pe = c2.paragraphs[0]
    rr = pe.add_run()
    rr.text = "Sites by prefecture"
    bp._style_run(rr, 16, bp.NAVY, bold=True, font=bp.HEAD_FONT)
    pj = c2.add_paragraph()
    rj = pj.add_run()
    rj.text = "県別スポット数"
    bp._style_run(rj, 11.5, bp.GREY, font=bp.JP_FONT)
    bp.en_jp(c2, f"Fukui {bp.POI_FUKUI} · Ishikawa {bp.POI_ISHIKAWA} · Toyama {bp.POI_TOYAMA}",
             f"福井 {bp.POI_FUKUI} ・ 石川 {bp.POI_ISHIKAWA} ・ 富山 {bp.POI_TOYAMA}",
             en_size=15, jp_size=12)
    bp.en_jp(c2, "Language means the review language, not the writer's nationality.",
             "「言語」は口コミの言語であり、書いた人の国籍ではない。",
             en_size=13, jp_size=10.5, en_color=bp.GREY, space_after=0)

    picture_fit(s, bp.PNG["volume"], Inches(6.55), Inches(2.0), Inches(6.2), Inches(3.5))
    bp.caption(s, Inches(6.65), Inches(5.6), Inches(6.0),
               "The rating model uses Japanese, English, and Chinese-language Google reviews.",
               "評価モデルでは、日本語・英語・中国語のGoogle口コミを使用。")
    bp.infobox(s, bp.MX, Inches(6.45), Inches(5.7), Inches(0.7),
               "We publish only counts and summaries. Raw review text stays inside the project.",
               "公開するのは件数と要約のみ。口コミの原文はプロジェクト内に保管する。")
    foot(s, 3)
    bp.notes(s,
             f"Lynn (45 sec): Our dataset is {bp.TAGGED_ROWS} Google reviews from {bp.N_POIS} sites in Fukui, "
             f"Ishikawa, and Toyama. The rating model uses {bp.LANG_JP} Japanese, {bp.LANG_EN} English, and "
             f"{bp.LANG_CN} Chinese-language reviews. 'Language' is the language of the text, not the reviewer's "
             "nationality. We only publish aggregate counts; raw text never leaves the project. "
             f"[For Q&A: a further {bp.LANG_OTHER} reviews were other languages and {bp.LANG_UNDETECTED} were "
             f"undetected or too short; a separate set of {bp.CN_ROWS} Xiaohongshu posts is shown later.]")

    # ===== SLIDE B (replaces live slide 4): Turning text into signals =========
    s = new()
    bp.label(s, "II.  METHODS")
    bp.title(s, "Turning review text into signals", "口コミを手がかりに変える")

    _, tf = bp.textbox(s, bp.MX, Inches(2.15), Inches(7.6), Inches(4.0))
    bp.en_jp(tf, f"1.  Tag each review with the {bp.N_ASPECTS} topics it mentions, using human-checked keyword lists.",
             f"1.  人が確認したキーワードを使い、各口コミが触れた{bp.N_ASPECTS}の話題に印をつける。",
             first=True, space_after=17)
    bp.en_jp(tf, f"2.  Call a review \"low-rated\" if it gave {bp.LOW_RATING_STARS} stars or fewer: {bp.LOW_RATING_ROWS} of {bp.MODEL_ROWS} reviews.",
             f"2.  {bp.LOW_RATING_STARS}つ星以下を「低評価」とする：{bp.MODEL_ROWS}件中{bp.LOW_RATING_ROWS}件。",
             space_after=17)
    bp.en_jp(tf, "3.  For each problem topic, ask: do reviews that mention it get low ratings more often?",
             "3.  各問題の話題について、それに触れた口コミは低評価が多いかを調べる。",
             space_after=17)
    bp.en_jp(tf, "The Google star rating is the shared yardstick. We do not compare sentiment-tool scores between languages.",
             "共通のものさしはGoogleの星評価。感情ツールの点数は言語間で比較しない。",
             bullet=True, space_after=0)

    bp.stat_callout(s, Inches(8.55), Inches(2.2), Inches(4.05), bp.MODEL_ROWS, "reviews modelled", "モデル対象の口コミ")
    bp.stat_callout(s, Inches(8.55), Inches(4.05), Inches(4.05), bp.LOW_RATING_ROWS, "were low-rated", "低評価だった")
    foot(s, 4)
    bp.notes(s,
             f"Andrew (50 sec): We tag each review with the {bp.N_ASPECTS} topics it mentions, using keyword lists a "
             f"person reviewed. A low rating means {bp.LOW_RATING_STARS} stars or fewer: {bp.LOW_RATING_ROWS} of "
             f"{bp.MODEL_ROWS} reviews. For each problem topic we ask whether mentioning it goes with a low rating "
             "more often. The Google star rating is our common outcome, so we never compare sentiment scores across "
             "different languages.")

    # ===== SLIDE C (merges live slides 5 + 6): How we tested it ===============
    s = new()
    bp.label(s, "II.  METHODS")
    bp.title(s, "How we tested it", "どう調べたか")

    _, tf = bp.textbox(s, bp.MX, Inches(2.0), Inches(6.5), Inches(4.6))
    bp.en_jp(tf, "For each topic we make a fair comparison: are reviews that mention the problem more likely to be low-rated, after allowing for review length, language, and prefecture?",
             "各話題で公平に比較する：口コミの長さ・言語・県の違いをそろえた上で、その問題に触れた口コミは低評価になりやすいか？",
             first=True, bullet=True, en_size=15, jp_size=11.5, space_after=15)
    bp.en_jp(tf, "We test many topics, so we tighten the bar to avoid being fooled by chance.",
             "多くの話題を調べるので、偶然の当たりに惑わされないよう基準を厳しくする。",
             bullet=True, en_size=15, jp_size=11.5, space_after=15)
    bp.en_jp(tf, f"A topic is kept only if the link is solid, points to worse ratings, and appears at least {bp.MIN_RANKING_MENTIONS} times.",
             f"残すのは、結びつきが確かで、悪い評価の方向で、{bp.MIN_RANKING_MENTIONS}回以上現れた話題だけ。",
             bullet=True, en_size=15, jp_size=11.5, space_after=15)
    bp.en_jp(tf, "We then sort each kept topic into \"better information can help\" or \"the site operator must act.\"",
             "残った話題を「情報で改善できる」か「事業者の対応が必要」かに振り分ける。",
             bullet=True, en_size=15, jp_size=11.5, space_after=0)

    picture_fit(s, FUNNEL_PNG, Inches(7.35), Inches(1.95), Inches(5.4), Inches(4.0))
    bp.infobox(s, bp.MX, Inches(6.5), Inches(12.1), Inches(0.66),
               "This ranks where to run the next experiment. It does not prove that any nudge works.",
               "これは「次にどこで実験するか」の順位づけ。ナッジの効果を証明するものではない。")
    foot(s, 5)
    bp.notes(s,
             "Andrew (60 sec spoken, plain): For each topic we run a fair comparison that accounts for review length, "
             "language, and prefecture. Because we test many topics, we tighten the bar to avoid false alarms. A topic "
             f"is kept only if the link is solid, harmful, and seen at least {bp.MIN_RANKING_MENTIONS} times. We then "
             "label each survivor as something information can help or something the operator must fix.\n\n"
             "--- FULL METHOD (for Q&A / examiners; do NOT read aloud) ---\n"
             f"Model: Firth bias-reduced (penalized) logistic regression, one per aspect; outcome = low rating "
             f"({bp.LOW_RATING_STARS} stars or fewer); adjusted for text length, review language, and prefecture. "
             "Firth was chosen because several aspects were rarely mentioned, which makes ordinary logistic regression "
             f"unstable. Output: an adjusted odds ratio with a 95% confidence interval and a p-value for each model "
             f"({bp.MODELS_PRIMARY} primary models, {bp.MODELS_TOTAL} across all segments). "
             f"Multiple testing: Benjamini-Hochberg FDR. Ranking gate: FDR-significant AND a harmful association "
             f"(odds ratio > 1) AND at least {bp.MIN_RANKING_MENTIONS} pooled mentions. Sanity check: for "
             f"{bp.FIRTH_SANITY_ASPECTS} selected aspects, Firth and standard-logit estimates agreed closely "
             f"(maximum |delta log OR| = {bp.FIRTH_SANITY}); smallest model {bp.MIN_P}. "
             "Caveat: POI-level clustering is not modeled in the Firth estimates, so row-level uncertainty may be "
             "understated; this is exploratory ranking, not a causal claim. Full table and CIs are in "
             "Seminar-Speaker-Script-and-QA.docx.")

    prs.save(str(OUT_PPTX))
    return prs


def main() -> int:
    build()
    size = OUT_PPTX.stat().st_size
    print(f"wrote {OUT_PPTX} ({size:,} bytes); 3 slides (replaces live methods slides 3-6)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
