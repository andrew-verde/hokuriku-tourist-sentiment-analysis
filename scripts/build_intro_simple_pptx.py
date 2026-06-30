#!/usr/bin/env python3
"""Build NUDGE-Intro-Simple.pptx: a plain, read-along replacement for the
INTRODUCTION, to be hand-inserted into the live Google Drive deck.

Splits the single overloaded "Research question" slide into TWO slides:
  Slide 2  - the actual research question (the WHAT), big and self-reading.
  Slide 3  - "What we tested" (the HOW), useful hypotheses + caveat.

Audience = Japanese viewers with limited English: every line is a full EN
sentence the presenter can read off, with a complete plain JP subtitle, so the
talk needs minimal extra explanation. Reuses scripts/build_nudge_pptx.py for
styling helpers, brand colors, and provenance-resolved constants (TAGGED_ROWS,
N_POIS). Nothing in the existing generation is modified.

Run:  .venv/bin/python3 scripts/build_intro_simple_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import build_nudge_pptx as bp

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "NUDGE-Intro-Simple.pptx"

# Splitting one intro slide into two grows the merged deck from 12 to 13.
DECK_TOTAL = 13


def foot(slide, num: int):
    tb, tf = bp.textbox(slide, Inches(11.0), Inches(7.04), Inches(1.83), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{num} / {DECK_TOTAL}"
    bp._style_run(r, 10, bp.GREY)


def _centered(tf, runs):
    """runs = list of (text, size, color, font, bold, space_after_pt)."""
    for i, (text, size, color, font, bold, sa) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(sa)
        r = p.add_run()
        r.text = text
        bp._style_run(r, size, color, bold=bold, font=font)


def build():
    prs = Presentation()
    prs.slide_width = bp.EMU_W
    prs.slide_height = bp.EMU_H
    blank = prs.slide_layouts[6]

    def new():
        s = prs.slides.add_slide(blank)
        bp._set_bg(s, bp.WHITE)
        return s

    # ===== SLIDE 2: the actual research question (WHAT) ======================
    s = new()
    bp.label(s, "I.  INTRODUCTION")
    bp.title(s, "Our research question", "私たちの研究の問い")

    _, q = bp.textbox(s, bp.MX, Inches(1.95), bp.CW, Inches(1.4))
    bp.en_jp(
        q,
        "What do tourist reviews tell us about how to improve Hokuriku tourism?",
        "口コミは、北陸の観光をどう良くできるかを教えてくれるか?",
        en_size=26, jp_size=16, en_color=bp.NAVY, bold=True, first=True, space_after=0,
    )

    subs = [
        ("Which problems lower ratings?", "どの問題が評価を下げるか?", "→ fix it  改善", bp.NAVY),
        ("Which good places get too few visitors?", "良い場所で、来訪が少ないのはどこか?",
         "→ promote it  推奨", bp.BLUE),
    ]
    cw, gap = Inches(6.0), Inches(0.33)
    xs = [bp.MX, bp.MX + cw + gap]
    for (en, jp, tag, accent), x in zip(subs, xs):
        _, c = bp.card(s, x, Inches(3.65), cw, Inches(1.95))
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _centered(c, [
            (en, 19, bp.NAVY, bp.HEAD_FONT, True, 2),
            (jp, 12.5, bp.GREY, bp.JP_FONT, False, 8),
            (tag, 15, accent, bp.HEAD_FONT, True, 0),
        ])

    _, b = bp.textbox(s, bp.MX, Inches(5.95), bp.CW, Inches(0.6))
    bp.en_jp(
        b,
        f"We read {bp.TAGGED_ROWS} reviews from {bp.N_POIS} Hokuriku sites to answer this.",
        f"北陸{bp.N_POIS}スポットの口コミ{bp.TAGGED_ROWS}件で、これに答える。",
        en_size=15, jp_size=11.5, en_color=bp.GREY, first=True, space_after=0,
    )
    bp.notes(s,
             "Andrew (20 sec): Our question is simple. What do tourist reviews tell us about how to improve "
             "Hokuriku tourism? It has two parts: which problems lower ratings, so we can fix them; and which "
             f"good places get too few visitors, so we can promote them. We read {bp.TAGGED_ROWS} reviews from "
             f"{bp.N_POIS} Hokuriku sites to answer this.")
    foot(s, 2)

    # ===== SLIDE 3: useful hypothesis tests (HOW) ============================
    s = new()
    bp.label(s, "I.  INTRODUCTION")
    bp.title(s, "What we tested", "私たちが検証したこと")

    _, tf = bp.textbox(s, bp.MX, Inches(2.15), Inches(7.4), Inches(3.7))
    bp.en_jp(tf, "We tested whether each review aspect was associated with a low rating.",
             "口コミの各要素が低評価と関連するかを検証した。",
             first=True, bullet=True, en_size=17, jp_size=12, space_after=16)
    bp.en_jp(tf, "We adjusted for review length, language, and prefecture.",
             "口コミの長さ、言語、県を調整した。",
             bullet=True, en_size=17, jp_size=12, space_after=16)
    bp.en_jp(tf, "We kept harmful associations that passed the FDR check.",
             "FDR検定を通過した、低評価との関連を残した。",
             bullet=True, en_size=17, jp_size=12, space_after=0)

    results = [
        ("These tests shaped the final ranking.", "この検定結果が最終順位を決めた。",
         "Opening hours, itinerary fit, and wayfinding passed the FDR check.",
         "営業時間、旅程適合、案内表示がFDR検定を通過した。"),
        ("Earlier language-gap tests moved to the appendix.", "以前の言語差検定は付録へ移した。",
         "Their measurement limits kept them from the headline.",
         "測定上の限界があるため、主要結果にはしなかった。"),
    ]
    for i, (en, jp, sub_en, sub_jp) in enumerate(results):
        _, c = bp.card(s, Inches(8.55), Inches(2.2) + Inches(i * 1.75), Inches(4.18), Inches(1.55))
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _centered(c, [
            (en, 14.5, bp.NAVY, bp.HEAD_FONT, True, 1),
            (jp, 11.5, bp.GREY, bp.JP_FONT, False, 6),
            (sub_en, 11.5, bp.INK, bp.EN_FONT, False, 1),
            (sub_jp, 10, bp.GREY, bp.JP_FONT, False, 0),
        ])

    bp.infobox(
        s, bp.MX, Inches(6.25), Inches(7.4), Inches(0.7),
        "Exploratory study: these are adjusted associations, not causal effects.",
        "探索的研究:これは調整済みの関連であり、因果効果ではない。",
    )
    bp.notes(s,
             "Andrew (30 sec): We began with hypotheses, not solutions. For each reviewed aspect, we tested "
             "whether its presence was associated with a low star rating. We adjusted for review length, "
             "language, and prefecture, then applied an FDR check. Opening hours, itinerary fit, and wayfinding "
             "were useful signals for the final ranking. Earlier language-gap tests moved to the appendix. "
             "These are exploratory associations, not causal effects.")
    foot(s, 3)

    prs.save(str(OUT_PPTX))
    return prs


def main() -> int:
    build()
    n = len(Presentation(str(OUT_PPTX)).slides)
    print(f"wrote {OUT_PPTX} ({OUT_PPTX.stat().st_size:,} bytes); {n} slides "
          f"(splits live slide 2 into research-question + approach; deck total -> {DECK_TOTAL})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
